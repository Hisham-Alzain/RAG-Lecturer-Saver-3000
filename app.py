# app.py
# Enhanced Multilingual RAG Chat - University Tutor
# With: True BM25 Hybrid Search, Arabic Morphology, HyDE, MMR, Evaluation, Context Compression

import streamlit as st
import hashlib
import uuid
import re
import math
import time
import json
import pickle
from pathlib import Path
from collections import defaultdict, OrderedDict
from typing import List, Dict, Tuple, Optional, Set
from dataclasses import dataclass, field
from datetime import datetime
import random

# Third-party imports
import numpy as np
from sentence_transformers import SentenceTransformer, CrossEncoder
from chromadb import PersistentClient
from chromadb.config import Settings
from pptx import Presentation
from pypdf import PdfReader
from openai import OpenAI
from rank_bm25 import BM25Okapi

# Arabic morphology - graceful fallback if not installed
try:
    from camel_tools.morphology.database import MorphologyDB
    from camel_tools.morphology.analyzer import Analyzer
    from camel_tools.disambig.mle import MLEDisambiguator
    from camel_tools.tokenizers.word import simple_word_tokenize

    CAMEL_AVAILABLE = True
except ImportError:
    CAMEL_AVAILABLE = False
    print(
        "⚠️ camel-tools not installed. Arabic morphology disabled. Install with: pip install camel-tools"
    )


# ============================================================================
# CONFIGURATION
# ============================================================================
@dataclass
class Config:
    """Central configuration for the RAG system."""

    # Chunking
    PARENT_CHUNK_SIZE: int = 1500
    PARENT_OVERLAP: int = 200
    CHILD_CHUNK_SIZE: int = 950
    CHILD_OVERLAP: int = 120

    # Retrieval
    TOP_K: int = 5
    MULTI_QUERY_COUNT: int = 3

    # BM25
    ENABLE_BM25: bool = True
    BM25_WEIGHT: float = 0.4  # Weight for BM25 in hybrid fusion (vector gets 1 - this)
    BM25_INDEX_PATH: str = "./bm25_index.pkl"

    # Reranking
    ENABLE_RERANKING: bool = True
    RERANK_MODEL: str = "cross-encoder/mmarco-mMiniLMv2-L12-H384-v1"
    RERANK_CANDIDATES: int = 30
    RERANK_TOP_K: int = 5

    # MMR (Maximal Marginal Relevance)
    ENABLE_MMR: bool = True
    MMR_LAMBDA: float = 0.7  # Balance: 1.0 = pure relevance, 0.0 = pure diversity
    MMR_TOP_K: int = 5

    # HyDE (Hypothetical Document Embedding)
    ENABLE_HYDE: bool = True

    # Context Compression
    ENABLE_CONTEXT_COMPRESSION: bool = True
    COMPRESSION_MAX_SENTENCES: int = 5

    # Arabic Morphology
    ENABLE_ARABIC_MORPHOLOGY: bool = True and CAMEL_AVAILABLE

    # Semantic Cache
    ENABLE_SEMANTIC_CACHE: bool = True
    CACHE_SIM_THRESHOLD: float = 0.92
    CACHE_MAX_ENTRIES: int = 200

    # CRAG (Corrective RAG)
    ENABLE_CRAG: bool = True
    CRAG_MIN_MAX_RELEVANCE: float = 0.35
    CRAG_MIN_MEAN_TOP3: float = 0.28
    CRAG_RETRIEVE_EXPANSION: int = 4

    # Evaluation
    EVAL_DATASET_PATH: str = "./eval_dataset.json"
    EVAL_RESULTS_PATH: str = "./eval_results.json"

    # Paths and Models
    CHROMA_DIR: str = "./chroma_db"
    EMBEDDING_MODEL: str = "intfloat/multilingual-e5-small"


CONFIG = Config()


# LLM Provider configurations
LLM_PROVIDERS = {
    "Groq (Llama 3.3 70B) ⚡": {
        "base_url": "https://api.groq.com/openai/v1",
        "model": "llama-3.3-70b-versatile",
        "name": "Groq",
        "get_key_url": "https://console.groq.com/keys",
        "notes": "Fastest inference, 14,400 req/day free",
    },
    "Groq (Llama 3.1 8B) 🆓": {
        "base_url": "https://api.groq.com/openai/v1",
        "model": "llama-3.1-8b-instant",
        "name": "Groq",
        "get_key_url": "https://console.groq.com/keys",
        "notes": "Fast & lightweight, good for testing",
    },
    "Cerebras (Llama 3.3 70B) 🚀": {
        "base_url": "https://api.cerebras.ai/v1",
        "model": "llama-3.3-70b",
        "name": "Cerebras",
        "get_key_url": "https://cloud.cerebras.ai/",
        "notes": "~1000 tok/sec, very fast",
    },
    "OpenRouter (Llama 3.1 8B) 🆓": {
        "base_url": "https://openrouter.ai/api/v1",
        "model": "meta-llama/llama-3.1-8b-instruct:free",
        "name": "OpenRouter",
        "get_key_url": "https://openrouter.ai/keys",
        "notes": "Free tier, requires OpenRouter key",
    },
    "OpenRouter (Mistral 7B) 🆓": {
        "base_url": "https://openrouter.ai/api/v1",
        "model": "mistralai/mistral-7b-instruct:free",
        "name": "OpenRouter",
        "get_key_url": "https://openrouter.ai/keys",
        "notes": "Free tier, fast responses",
    },
}


# ============================================================================
# PAGE SETUP & STYLING
# ============================================================================
st.set_page_config(
    page_title="RAG Chat",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# CSS Styles
CUSTOM_CSS = """
<style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    
    .main .block-container {
        max-width: 900px;
        padding-top: 2rem;
        padding-bottom: 2rem;
    }
    
    .user-message {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1rem 1.25rem;
        border-radius: 1.25rem;
        margin: 1rem 0;
        max-width: 80%;
        margin-left: auto;
        font-size: 1rem;
    }
    
    .assistant-message {
        background: #f7f7f8;
        color: #1a1a1a;
        padding: 1.5rem;
        border-radius: 1rem;
        margin: 1rem 0;
        line-height: 1.7;
        font-size: 1rem;
    }
    
    .assistant-message-rtl {
        background: #f7f7f8;
        color: #1a1a1a;
        padding: 1.5rem;
        border-radius: 1rem;
        margin: 1rem 0;
        line-height: 1.9;
        font-size: 1.05rem;
        direction: rtl;
        text-align: right;
    }
    
    .source-card {
        background: white;
        border: 1px solid #e5e5e5;
        border-radius: 0.75rem;
        padding: 0.75rem 1rem;
        font-size: 0.85rem;
        display: flex;
        align-items: center;
        gap: 0.5rem;
        transition: all 0.2s;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
        margin-bottom: 0.5rem;
    }
    
    .source-card:hover {
        border-color: #667eea;
        box-shadow: 0 2px 8px rgba(102,126,234,0.15);
    }
    
    .source-number {
        background: #667eea;
        color: white;
        min-width: 1.5rem;
        height: 1.5rem;
        border-radius: 50%;
        display: flex;
        align-items: center;
        justify-content: center;
        font-size: 0.75rem;
        font-weight: 600;
    }
    
    .citation {
        background: #667eea;
        color: white;
        padding: 0.1rem 0.4rem;
        border-radius: 0.25rem;
        font-size: 0.75rem;
        font-weight: 600;
        margin: 0 0.1rem;
    }
    
    .section-label {
        color: #888;
        font-size: 0.8rem;
        font-weight: 600;
        text-transform: uppercase;
        letter-spacing: 0.05em;
        margin-bottom: 0.5rem;
    }
    
    .app-header {
        text-align: center;
        padding: 2rem 0;
    }
    
    .app-title {
        font-size: 2rem;
        font-weight: 700;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.5rem;
    }
    
    .app-subtitle {
        color: #666;
        font-size: 1rem;
    }
    
    .provider-badge {
        background: #e8f4e8;
        color: #2d6a2d;
        padding: 0.25rem 0.5rem;
        border-radius: 0.25rem;
        font-size: 0.75rem;
        font-weight: 500;
    }
    
    .eval-metric {
        background: #f0f4ff;
        border: 1px solid #d0d8ff;
        border-radius: 0.5rem;
        padding: 0.5rem 1rem;
        margin: 0.25rem;
        display: inline-block;
    }
    
    .eval-metric-name {
        font-size: 0.75rem;
        color: #666;
        text-transform: uppercase;
    }
    
    .eval-metric-value {
        font-size: 1.25rem;
        font-weight: 600;
        color: #333;
    }
</style>
"""
st.markdown(CUSTOM_CSS, unsafe_allow_html=True)


# ============================================================================
# SESSION STATE INITIALIZATION
# ============================================================================
def init_session_state():
    """Initialize all session state variables."""
    defaults = {
        "chat_history": [],
        "api_key": "",
        "selected_provider": list(LLM_PROVIDERS.keys())[0],
        "semantic_cache": OrderedDict(),
        "bm25_index": None,
        "bm25_corpus": [],
        "bm25_doc_ids": [],
        "arabic_analyzer": None,
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value


init_session_state()


# ============================================================================
# ARABIC MORPHOLOGY PROCESSOR
# ============================================================================
class ArabicMorphologyProcessor:
    """
    Handles Arabic text normalization and lemmatization using CAMeL Tools.

    Arabic is morphologically complex - words have roots, patterns, prefixes, and suffixes.
    For example: يكتبون (they write), كتاب (book), مكتوب (written) all share root ك-ت-ب.
    Lemmatization helps match these related forms during retrieval.
    """

    def __init__(self):
        self.analyzer = None
        self.disambiguator = None
        self._initialized = False

    def initialize(self):
        """Lazy initialization of CAMeL tools (they're heavy to load)."""
        if self._initialized or not CAMEL_AVAILABLE:
            return

        try:
            # Load the morphology database and analyzer
            db = MorphologyDB.builtin_db()
            self.analyzer = Analyzer(db)
            self.disambiguator = MLEDisambiguator.pretrained()
            self._initialized = True
        except Exception as e:
            print(f"Failed to initialize Arabic morphology: {e}")
            self._initialized = False

    def normalize_arabic(self, text: str) -> str:
        """
        Normalize Arabic text for consistent matching.

        This handles common variations in Arabic writing:
        - Alef variants (أ إ آ ا) -> ا
        - Teh marbuta (ة) -> ه
        - Tatweel/kashida (ـ) removal
        - Diacritics (tashkeel) removal
        """
        if not text:
            return text

        # Normalize Alef variants
        text = re.sub(r"[أإآٱ]", "ا", text)

        # Normalize Alef Maksura
        text = re.sub(r"ى", "ي", text)

        # Normalize Teh Marbuta (optional - sometimes you want to keep it)
        # text = re.sub(r'ة', 'ه', text)

        # Remove Tatweel (elongation character)
        text = re.sub(r"ـ", "", text)

        # Remove diacritics (tashkeel)
        arabic_diacritics = re.compile(r"[\u064B-\u065F\u0670]")
        text = arabic_diacritics.sub("", text)

        return text

    def lemmatize(self, text: str) -> str:
        """
        Lemmatize Arabic text to get base forms.

        This converts inflected forms to their lemmas (dictionary forms),
        which dramatically improves recall for Arabic queries.
        """
        if not CAMEL_AVAILABLE or not self._initialized:
            return self.normalize_arabic(text)

        try:
            # Tokenize the text
            tokens = simple_word_tokenize(text)

            # Disambiguate and get lemmas
            disambiguated = self.disambiguator.disambiguate(tokens)

            lemmas = []
            for d in disambiguated:
                # Get the best analysis
                if d.analyses:
                    # Use the lemma from the top analysis
                    lemma = d.analyses[0].analysis.get("lex", d.word)
                    lemmas.append(lemma)
                else:
                    lemmas.append(d.word)

            return " ".join(lemmas)
        except Exception as e:
            # Fallback to normalization only
            return self.normalize_arabic(text)

    def tokenize_for_bm25(self, text: str) -> List[str]:
        """
        Tokenize and lemmatize text for BM25 indexing.

        Returns a list of normalized/lemmatized tokens suitable for BM25.
        """
        lang = LanguageUtils.detect_language(text)

        if lang == "ar":
            # Normalize first
            normalized = self.normalize_arabic(text)

            # Lemmatize if available
            if self._initialized:
                processed = self.lemmatize(normalized)
            else:
                processed = normalized

            # Tokenize
            if CAMEL_AVAILABLE:
                tokens = simple_word_tokenize(processed)
            else:
                tokens = re.findall(r"[\u0600-\u06ff]+", processed)

            # Remove stop words
            stop_words = LanguageUtils.ARABIC_STOP_WORDS
            tokens = [t for t in tokens if t not in stop_words and len(t) > 1]
        else:
            # English tokenization
            tokens = re.findall(r"\b\w+\b", text.lower())
            stop_words = LanguageUtils.ENGLISH_STOP_WORDS
            tokens = [t for t in tokens if t not in stop_words and len(t) > 2]

        return tokens


# Global Arabic processor instance
_arabic_processor = None


def get_arabic_processor() -> ArabicMorphologyProcessor:
    """Get or create the Arabic morphology processor."""
    global _arabic_processor
    if _arabic_processor is None:
        _arabic_processor = ArabicMorphologyProcessor()
        if CONFIG.ENABLE_ARABIC_MORPHOLOGY:
            _arabic_processor.initialize()
    return _arabic_processor


# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================
class LanguageUtils:
    """Utilities for language detection and handling."""

    # Arabic character range
    ARABIC_RANGE = ("\u0600", "\u06ff")

    # Arabic stop words for keyword extraction
    ARABIC_STOP_WORDS = {
        "من",
        "في",
        "على",
        "إلى",
        "عن",
        "مع",
        "هذا",
        "هذه",
        "ذلك",
        "تلك",
        "التي",
        "الذي",
        "هو",
        "هي",
        "هم",
        "هن",
        "أن",
        "إن",
        "كان",
        "كانت",
        "يكون",
        "تكون",
        "ما",
        "لا",
        "لم",
        "لن",
        "قد",
        "كل",
        "بعض",
        "أي",
        "و",
        "أو",
        "ثم",
        "حتى",
        "إذا",
        "لو",
        "كما",
        "بل",
        "لكن",
        "غير",
        "بين",
        "عند",
        "منذ",
        "حول",
        "خلال",
        "بعد",
        "قبل",
        "فوق",
        "تحت",
        "ال",
        "الى",
        "عن",
        "هل",
        "كيف",
        "لماذا",
        "متى",
        "أين",
        "من",
    }

    # English stop words
    ENGLISH_STOP_WORDS = {
        "the",
        "a",
        "an",
        "and",
        "or",
        "but",
        "in",
        "on",
        "at",
        "to",
        "for",
        "of",
        "with",
        "is",
        "are",
        "was",
        "were",
        "be",
        "been",
        "being",
        "have",
        "has",
        "had",
        "do",
        "does",
        "did",
        "what",
        "which",
        "who",
        "when",
        "where",
        "why",
        "how",
        "this",
        "that",
        "these",
        "those",
        "it",
        "its",
        "as",
        "by",
        "from",
        "about",
        "into",
        "through",
    }

    @classmethod
    def detect_language(cls, text: str) -> str:
        """Detect if text is primarily Arabic or English."""
        if not text:
            return "en"

        arabic_chars = sum(
            1 for c in text if cls.ARABIC_RANGE[0] <= c <= cls.ARABIC_RANGE[1]
        )
        total_alpha = sum(1 for c in text if c.isalpha())

        if total_alpha == 0:
            return "en"

        return "ar" if (arabic_chars / total_alpha) > 0.3 else "en"

    @classmethod
    def get_stop_words(cls, lang: str) -> set:
        """Get stop words for the specified language."""
        return cls.ARABIC_STOP_WORDS if lang == "ar" else cls.ENGLISH_STOP_WORDS

    @classmethod
    def get_text_separators(cls, lang: str) -> List[str]:
        """Get text separators for chunking based on language."""
        if lang == "ar":
            return ["\n## ", "\n### ", "\n\n", "。", "\n", ". ", "، ", " ", ""]
        return ["\n## ", "\n### ", "\n\n", ". ", "\n", " ", ""]


class MathUtils:
    """Mathematical utility functions."""

    @staticmethod
    def sigmoid(x: float) -> float:
        """Numerically stable sigmoid function."""
        if x >= 0:
            z = math.exp(-x)
            return 1 / (1 + z)
        z = math.exp(x)
        return z / (1 + z)

    @staticmethod
    def cosine_sim(a: List[float], b: List[float]) -> float:
        """Compute cosine similarity (assumes normalized vectors)."""
        return float(sum(x * y for x, y in zip(a, b)))

    @staticmethod
    def cosine_sim_matrix(embeddings: np.ndarray) -> np.ndarray:
        """Compute pairwise cosine similarity matrix."""
        # Normalize embeddings
        norms = np.linalg.norm(embeddings, axis=1, keepdims=True)
        normalized = embeddings / (norms + 1e-10)
        # Compute similarity matrix
        return np.dot(normalized, normalized.T)


# ============================================================================
# LLM CLIENT
# ============================================================================
class LLMClient:
    """Wrapper for LLM API interactions with proper provider-specific handling."""

    def __init__(self, provider_key: str, api_key: str):
        config = LLM_PROVIDERS[provider_key]
        self.model = config["model"]
        self.provider_name = config["name"]
        self.provider_key = provider_key

        # Build default headers - OpenRouter requires these additional headers
        default_headers = {}
        if "openrouter" in config["base_url"].lower():
            default_headers = {
                "HTTP-Referer": "https://university-tutor-rag.streamlit.app",  # Required by OpenRouter
                "X-Title": "University Tutor RAG",  # Optional but recommended
            }

        self.client = OpenAI(
            base_url=config["base_url"],
            api_key=api_key,
            default_headers=default_headers if default_headers else None,
        )

    def test_connection(self) -> Tuple[bool, str]:
        """
        Test if the API connection works.
        Returns (success: bool, message: str)
        """
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[{"role": "user", "content": "Hi"}],
                max_tokens=5,
            )
            return True, "Connection successful"
        except Exception as e:
            error_msg = str(e)

            # Provide helpful error messages based on common issues
            if "403" in error_msg:
                if "openrouter" in self.provider_key.lower():
                    return (
                        False,
                        "403 Error: Check that your OpenRouter API key is valid and has credits.",
                    )
                elif "groq" in self.provider_key.lower():
                    return (
                        False,
                        "403 Error: Check that your Groq API key is valid. Get one at console.groq.com/keys",
                    )
                elif "cerebras" in self.provider_key.lower():
                    return (
                        False,
                        "403 Error: Check that your Cerebras API key is valid.",
                    )
                else:
                    return (
                        False,
                        f"403 Error: Access denied. Verify your API key matches the selected provider ({self.provider_name}).",
                    )
            elif "401" in error_msg:
                return False, "401 Error: Invalid API key. Please check your key."
            elif "429" in error_msg:
                return False, "429 Error: Rate limited. Wait a moment and try again."
            elif "model" in error_msg.lower() and "not found" in error_msg.lower():
                return (
                    False,
                    f"Model '{self.model}' not found. The model may have been renamed or is not available.",
                )
            else:
                return False, f"Connection error: {error_msg}"

    def chat(
        self, messages: List[Dict], temperature: float = 0.7, max_tokens: int = 2048
    ) -> str:
        """Send chat completion request with improved error handling."""
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=messages,
                temperature=temperature,
                max_tokens=max_tokens,
            )
            return response.choices[0].message.content
        except Exception as e:
            error_msg = str(e)

            # Provide more helpful error messages
            if "403" in error_msg:
                return f"Error: Access denied (403). Please verify:\n1. Your API key is correct\n2. The key matches the selected provider ({self.provider_name})\n3. Your account has available credits/quota"
            elif "401" in error_msg:
                return "Error: Invalid API key (401). Please check your API key."
            elif "429" in error_msg:
                return "Error: Rate limited (429). Please wait a moment and try again."
            elif "model" in error_msg.lower():
                return f"Error: Model issue - {error_msg}. Try selecting a different provider."
            else:
                return f"Error: {error_msg}"

    def translate_query(self, query: str, source_lang: str) -> Optional[str]:
        """Translate query to the other language for bilingual search."""
        if source_lang == "ar":
            instruction = (
                "Translate to English. Output ONLY the translation, nothing else:"
            )
        else:
            instruction = (
                "Translate to Arabic. Output ONLY the translation, nothing else:"
            )

        try:
            result = self.chat(
                [{"role": "user", "content": f"{instruction}\n{query}"}],
                temperature=0.3,
            )
            if (
                result
                and not result.startswith("Error:")
                and result.lower() != query.lower()
            ):
                return result.strip()
        except Exception:
            pass
        return None

    def generate_multi_queries(
        self, original_query: str, question_lang: str
    ) -> List[str]:
        """Generate multiple query variations for better retrieval coverage."""
        if question_lang == "ar":
            instruction = """أنشئ 3 صيغ بديلة لهذا السؤال بالعربية.
ركز على:
1. مصطلحات تقنية أكثر تحديداً
2. صيغة أوسع سياقاً
3. صياغة مختلفة بمرادفات

اكتب الأسئلة فقط، سؤال في كل سطر، بدون ترقيم."""
        else:
            instruction = """Generate 3 alternative ways to ask this question in English.
Focus on:
1. More specific technical terms
2. Broader context version
3. Different phrasing with synonyms

Output ONLY the questions, one per line, without numbering."""

        try:
            response = self.chat(
                [
                    {
                        "role": "user",
                        "content": f"{instruction}\n\nOriginal Question: {original_query}",
                    }
                ],
                temperature=0.8,
            )

            lines = [l.strip() for l in response.strip().split("\n") if l.strip()]
            cleaned = []
            for line in lines:
                line = re.sub(r"^[\d\.\-\*\)]+\s*", "", line)
                if line and len(line) > 10:
                    cleaned.append(line)

            return cleaned[: CONFIG.MULTI_QUERY_COUNT]
        except Exception as e:
            st.warning(f"Multi-query generation failed: {e}")
            return []

    def generate_hypothetical_answer(
        self, query: str, question_lang: str
    ) -> Optional[str]:
        """
        Generate a hypothetical answer for HyDE (Hypothetical Document Embedding).

        HyDE works by generating what a good answer WOULD look like, then embedding
        that hypothetical answer to search for similar real documents. This bridges
        the gap between question-style text and document-style text.
        """
        if question_lang == "ar":
            instruction = """أنت مساعد أكاديمي. اكتب فقرة قصيرة (3-4 جمل) تجيب على هذا السؤال كما لو كانت من مادة محاضرة.
اكتب بأسلوب أكاديمي موضوعي. لا تقل "السؤال هو" أو "الجواب هو".
اكتب المحتوى مباشرة كما يظهر في كتاب أو محاضرة."""
        else:
            instruction = """You are an academic assistant. Write a short paragraph (3-4 sentences) answering this question as if from lecture material.
Write in objective academic style. Don't say "The question is" or "The answer is".
Write the content directly as it would appear in a textbook or lecture."""

        try:
            result = self.chat(
                [{"role": "user", "content": f"{instruction}\n\nQuestion: {query}"}],
                temperature=0.5,
                max_tokens=256,
            )
            if result and not result.startswith("Error:"):
                return result.strip()
        except Exception:
            pass
        return None

    def compress_context(
        self, query: str, context: str, question_lang: str, max_sentences: int = 5
    ) -> str:
        """
        Compress context by extracting only the most relevant sentences.

        This reduces noise in the context sent to the generator, improving
        answer quality and reducing hallucination.
        """
        if question_lang == "ar":
            instruction = f"""استخرج أهم {max_sentences} جمل من النص التالي التي تجيب مباشرة على السؤال.
اكتب الجمل فقط، كل جملة في سطر منفصل.
إذا لم تجد معلومات ذات صلة، اكتب "لا توجد معلومات ذات صلة"."""
        else:
            instruction = f"""Extract the {max_sentences} most important sentences from the following text that directly answer the question.
Write only the sentences, one per line.
If no relevant information found, write "No relevant information found"."""

        try:
            result = self.chat(
                [
                    {
                        "role": "user",
                        "content": f"{instruction}\n\nQuestion: {query}\n\nText:\n{context}",
                    }
                ],
                temperature=0.3,
                max_tokens=512,
            )
            if result and not result.startswith("Error:"):
                return result.strip()
        except Exception:
            pass
        return context  # Return original if compression fails


# ============================================================================
# MODEL LOADERS (Cached)
# ============================================================================
@st.cache_resource
def load_embedder() -> SentenceTransformer:
    """Load the embedding model."""
    return SentenceTransformer(CONFIG.EMBEDDING_MODEL)


@st.cache_resource
def load_reranker() -> CrossEncoder:
    """Load the cross-encoder reranker model."""
    return CrossEncoder(CONFIG.RERANK_MODEL)


def get_chroma_collection():
    """Get or create ChromaDB collection (cosine space for normalized embeddings)."""
    if "chroma_collection" not in st.session_state:
        chroma_client = PersistentClient(
            path=CONFIG.CHROMA_DIR, settings=Settings(allow_reset=True)
        )
        st.session_state["chroma_client"] = chroma_client

        # IMPORTANT: cosine distance (works with normalize_embeddings=True)
        st.session_state["chroma_collection"] = chroma_client.get_or_create_collection(
            "lecture_rag",
            metadata={"hnsw:space": "cosine"},
        )
    return st.session_state["chroma_client"], st.session_state["chroma_collection"]


# ============================================================================
# BM25 INDEX MANAGER
# ============================================================================
class BM25IndexManager:
    """
    Manages the BM25 sparse retrieval index.

    BM25 (Best Match 25) is a bag-of-words retrieval function that ranks documents
    based on term frequency and inverse document frequency. It excels at exact
    keyword matching, which complements vector search's semantic matching.

    This is crucial for Arabic because:
    1. Technical terms and proper nouns often don't embed well
    2. Exact matches for transliterated terms (e.g., "CNN", "BERT") are important
    3. Arabic morphological variants can be pre-processed for better matching
    """

    def __init__(self):
        self.index: Optional[BM25Okapi] = None
        self.corpus: List[List[str]] = []  # Tokenized documents
        self.doc_ids: List[str] = []  # Corresponding document IDs
        self.raw_texts: List[str] = []  # Original texts for reference
        self.arabic_processor = get_arabic_processor()

    def tokenize_document(self, text: str) -> List[str]:
        """Tokenize a document for BM25 indexing with Arabic morphology support."""
        return self.arabic_processor.tokenize_for_bm25(text)

    def build_index(self, documents: List[Dict]):
        """
        Build BM25 index from a list of documents.

        Each document should have 'id' and 'text' keys.
        """
        self.corpus = []
        self.doc_ids = []
        self.raw_texts = []

        for doc in documents:
            doc_id = doc.get("id", str(uuid.uuid4()))
            text = doc.get("text", "")

            tokens = self.tokenize_document(text)

            if tokens:  # Only add if we have tokens
                self.corpus.append(tokens)
                self.doc_ids.append(doc_id)
                self.raw_texts.append(text)

        if self.corpus:
            self.index = BM25Okapi(self.corpus)

    def add_documents(self, documents: List[Dict]):
        """Add new documents to the existing index."""
        for doc in documents:
            doc_id = doc.get("id", str(uuid.uuid4()))
            text = doc.get("text", "")

            tokens = self.tokenize_document(text)

            if tokens:
                self.corpus.append(tokens)
                self.doc_ids.append(doc_id)
                self.raw_texts.append(text)

        # Rebuild index with new documents
        if self.corpus:
            self.index = BM25Okapi(self.corpus)

    def search(self, query: str, top_k: int = 10) -> List[Tuple[str, float]]:
        """
        Search the BM25 index and return top-k results.

        Returns list of (doc_id, score) tuples.
        """
        if not self.index or not self.corpus:
            return []

        # Tokenize query with same preprocessing as documents
        query_tokens = self.tokenize_document(query)

        if not query_tokens:
            return []

        # Get BM25 scores for all documents
        scores = self.index.get_scores(query_tokens)

        # Get top-k indices
        top_indices = np.argsort(scores)[::-1][:top_k]

        results = []
        for idx in top_indices:
            if scores[idx] > 0:  # Only include documents with positive scores
                results.append((self.doc_ids[idx], float(scores[idx])))

        return results

    def save(self, path: str):
        """Save the BM25 index to disk."""
        data = {
            "corpus": self.corpus,
            "doc_ids": self.doc_ids,
            "raw_texts": self.raw_texts,
        }
        with open(path, "wb") as f:
            pickle.dump(data, f)

    def load(self, path: str) -> bool:
        """Load the BM25 index from disk."""
        try:
            with open(path, "rb") as f:
                data = pickle.load(f)

            self.corpus = data["corpus"]
            self.doc_ids = data["doc_ids"]
            self.raw_texts = data.get("raw_texts", [])

            if self.corpus:
                self.index = BM25Okapi(self.corpus)
            return True
        except Exception as e:
            print(f"Failed to load BM25 index: {e}")
            return False


# Global BM25 manager
def get_bm25_manager() -> BM25IndexManager:
    """Get or create the BM25 index manager."""
    if "bm25_manager" not in st.session_state:
        manager = BM25IndexManager()
        # Try to load existing index
        if Path(CONFIG.BM25_INDEX_PATH).exists():
            manager.load(CONFIG.BM25_INDEX_PATH)
        st.session_state["bm25_manager"] = manager
    return st.session_state["bm25_manager"]


# ============================================================================
# DOCUMENT PROCESSING
# ============================================================================
class DocumentProcessor:
    """Handle document extraction and chunking."""

    def __init__(self, embedder: SentenceTransformer):
        self.embedder = embedder

    def get_token_count(self, text: str) -> int:
        """Get actual token count using the embedding model's tokenizer."""
        try:
            tokenizer = self.embedder.tokenizer
            tokens = tokenizer.encode(text, add_special_tokens=True)
            return len(tokens)
        except Exception:
            return int(len(text.split()) * 1.3)

    def extract_text_from_pdf(self, file) -> List[Tuple[int, str]]:
        """Extract text from PDF file."""
        try:
            reader = PdfReader(file)
            pages = []
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append((page_num, text))
            return pages
        except Exception as e:
            st.error(f"PDF error: {e}")
            return []

    def extract_text_from_ppt(self, file) -> List[Tuple[int, str]]:
        """Extract text from PowerPoint file."""
        try:
            prs = Presentation(file)
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    slides.append((idx, "\n".join(slide_text)))
            return slides
        except Exception as e:
            st.error(f"PPT error: {e}")
            return []

    def split_text_recursive(
        self, text: str, max_size: int, overlap: int, separators: List[str]
    ) -> List[str]:
        """
        Character-based recursive splitter.

        max_size/overlap are in CHARACTERS (matches homework "500-1000 chars").
        """
        text = (text or "").strip()
        if not text:
            return []

        if len(text) <= max_size:
            return [text]

        for sep in separators:
            if sep and sep in text:
                parts = text.split(sep)
                chunks = []
                current = ""

                for i, part in enumerate(parts):
                    if i > 0:
                        part = sep + part

                    candidate = current + part
                    if len(candidate) <= max_size:
                        current = candidate
                    else:
                        if current:
                            chunks.append(current.strip())

                        # overlap in characters
                        if overlap > 0 and current:
                            tail = current[-overlap:]
                            current = (tail + part).strip()
                        else:
                            current = part.strip()

                if current:
                    chunks.append(current.strip())

                # If split worked into >1 chunk, stop here
                if len(chunks) > 1:
                    return chunks

        # Fallback: sliding window char split
        step = max(1, max_size - overlap)
        out = []
        for start in range(0, len(text), step):
            out.append(text[start : start + max_size].strip())
            if start + max_size >= len(text):
                break
        return [c for c in out if c]


    def hierarchical_chunk_text(self, text: str, metadata: Dict) -> List[Dict]:
        """Create hierarchical parent-child chunks."""
        if not text or not text.strip():
            return []

        lang = LanguageUtils.detect_language(text)
        separators = LanguageUtils.get_text_separators(lang)

        # Create parent chunks
        parent_chunks = self.split_text_recursive(
            text, CONFIG.PARENT_CHUNK_SIZE, CONFIG.PARENT_OVERLAP, separators
        )

        all_chunks = []
        for parent_idx, parent_text in enumerate(parent_chunks):
            # Create child chunks from each parent
            child_chunks = self.split_text_recursive(
                parent_text, CONFIG.CHILD_CHUNK_SIZE, CONFIG.CHILD_OVERLAP, separators
            )

            for child_idx, child_text in enumerate(child_chunks):
                # Add context prefix for better embeddings
                context_prefix = f"[{metadata['source']} | page {metadata['page']}]\n"


                child_with_context = context_prefix + child_text

                all_chunks.append(
                    {
                        "child_text": child_with_context,
                        "child_text_raw": child_text,
                        "parent_text": parent_text,
                        "parent_id": f"{metadata['source']}_p{metadata['page']}_parent{parent_idx}",
                        "metadata": {
                            **metadata,
                            "parent_idx": parent_idx,
                            "child_idx": child_idx,
                            "chunk_type": "child",
                        },
                    }
                )

        return all_chunks


class DocumentIngester:
    """Handle document ingestion into both vector database and BM25 index."""

    def __init__(self, embedder: SentenceTransformer, collection):
        self.embedder = embedder
        self.collection = collection
        self.processor = DocumentProcessor(embedder)
        self.bm25_manager = get_bm25_manager()

    @staticmethod
    def get_file_hash(file) -> str:
        """Generate MD5 hash of file contents."""
        file.seek(0)
        file_hash = hashlib.md5(file.read()).hexdigest()
        file.seek(0)
        return file_hash

    def file_already_ingested(self, file_hash: str) -> bool:
        """Check if file has already been ingested."""
        try:
            results = self.collection.get(where={"file_hash": file_hash}, limit=1)
            return len(results["ids"]) > 0
        except Exception:
            return False

    def ingest(self, file) -> Tuple[int, str]:
        """Ingest a document into both vector database and BM25 index."""
        file_hash = self.get_file_hash(file)

        if self.file_already_ingested(file_hash):
            return 0, "exists"

        filename = file.name.lower()
        if filename.endswith(".pdf"):
            pages = self.processor.extract_text_from_pdf(file)
            source_type = "pdf"
        elif filename.endswith((".ppt", ".pptx")):
            pages = self.processor.extract_text_from_ppt(file)
            source_type = "ppt"
        else:
            return 0, "unsupported"

        if not pages:
            return 0, "no text"

        all_chunks = []
        chunk_metadata = []
        bm25_docs = []  # Documents for BM25 indexing

        for page_num, text in pages:
            chunk_lang = LanguageUtils.detect_language(text)

            chunks = self.processor.hierarchical_chunk_text(
                text,
                {
                    "source": file.name,
                    "page": page_num,
                    "type": source_type,
                    "file_hash": file_hash,
                    "lang": chunk_lang,
                },
            )

            for chunk_data in chunks:
                base = f"{file_hash}|{file.name}|p{page_num}|pa{chunk_data['metadata']['parent_idx']}|ch{chunk_data['metadata']['child_idx']}"
                chunk_id = hashlib.md5(base.encode("utf-8")).hexdigest()
                all_chunks.append(f"passage: {chunk_data['child_text']}")
                chunk_metadata.append(
                    {
                        "id": chunk_id,
                        "child_text": chunk_data["child_text"],
                        "child_text_raw": chunk_data["child_text_raw"],
                        "parent_text": chunk_data["parent_text"],
                        "parent_id": chunk_data["parent_id"],
                        **chunk_data["metadata"],
                    }
                )

                # Prepare document for BM25 - use parent text for richer context
                bm25_docs.append(
                    {
                        "id": chunk_id,
                        "text": chunk_data["child_text_raw"],
                    }
                )

        if not all_chunks:
            return 0, "no chunks"

        try:
            # Generate embeddings for vector store
            embeddings = self.embedder.encode(
                all_chunks, normalize_embeddings=True, show_progress_bar=False
            ).tolist()

            # Add to ChromaDB
            self.collection.add(
                ids=[m["id"] for m in chunk_metadata],
                documents=[m["child_text_raw"] for m in chunk_metadata],
                embeddings=embeddings,
                metadatas=[
                    {
                        "source": m["source"],
                        "page": m["page"],
                        "type": m["type"],
                        "file_hash": m["file_hash"],
                        "lang": m["lang"],
                        "parent_text": m["parent_text"],
                        "parent_id": m["parent_id"],
                        "parent_idx": m.get("parent_idx", 0),
                        "child_idx": m.get("child_idx", 0),
                        "chunk_type": m.get("chunk_type", "child"),
                    }
                    for m in chunk_metadata
                ],
            )

            # Add to BM25 index
            if CONFIG.ENABLE_BM25:
                self.bm25_manager.add_documents(bm25_docs)
                self.bm25_manager.save(CONFIG.BM25_INDEX_PATH)

            return len(all_chunks), "success"
        except Exception as e:
            return 0, str(e)


# ============================================================================
# RETRIEVAL SYSTEM
# ============================================================================
class Retriever:
    """
    Handle document retrieval with true hybrid search, HyDE, reranking, and MMR.

    The retrieval pipeline:
    1. Multi-query generation (semantic + keyword variations)
    2. Cross-lingual translation (Arabic <-> English)
    3. HyDE (Hypothetical Document Embedding) for abstract queries
    4. Hybrid search (Vector + BM25 with RRF fusion)
    5. Cross-encoder reranking
    6. MMR diversity filtering
    7. CRAG quality assessment
    """

    def __init__(
        self, embedder: SentenceTransformer, collection, llm_client: LLMClient
    ):
        self.embedder = embedder
        self.collection = collection
        self.llm_client = llm_client
        self.reranker = None
        self.bm25_manager = get_bm25_manager()

        if CONFIG.ENABLE_RERANKING or CONFIG.ENABLE_CRAG:
            try:
                self.reranker = load_reranker()
            except Exception as e:
                st.warning(f"Reranker failed to load: {e}")

    def extract_keywords(self, text: str, max_keywords: int = 5) -> List[str]:
        """Extract keywords from text with Arabic morphology support."""
        arabic_processor = get_arabic_processor()
        tokens = arabic_processor.tokenize_for_bm25(text)

        # Return unique keywords
        seen = set()
        unique_keywords = []
        for kw in tokens:
            if kw not in seen:
                seen.add(kw)
                unique_keywords.append(kw)
            if len(unique_keywords) >= max_keywords:
                break

        return unique_keywords

    def vector_search(self, query: str, top_k: int = 10) -> List[Dict]:
        """Perform vector similarity search."""
        query_embedding = self.embedder.encode(
            [f"query: {query}"], normalize_embeddings=True
        ).tolist()

        results = self.collection.query(
            query_embeddings=query_embedding, n_results=top_k
        )

        if not results["ids"] or not results["ids"][0]:
            return []

        return [
            {
                "id": doc_id,
                "score": max(0.0, 1.0 - float(dist)),  # Convert distance to similarity
                "metadata": meta,
                "document": doc,
                "distance": dist,
            }
            for doc_id, dist, meta, doc in zip(
                results["ids"][0],
                results["distances"][0],
                results["metadatas"][0],
                results["documents"][0],
            )
        ]

    def bm25_search(self, query: str, top_k: int = 10) -> List[Dict]:
        """
        Perform BM25 sparse retrieval search.

        This catches exact term matches that vector search might miss,
        especially important for:
        - Technical terms and acronyms (CNN, BERT, API)
        - Arabic proper nouns and transliterations
        - Exact keyword matches
        """
        if not CONFIG.ENABLE_BM25 or not self.bm25_manager.index:
            return []

        bm25_results = self.bm25_manager.search(query, top_k=top_k)

        if not bm25_results:
            return []

        # Fetch full document info from ChromaDB
        doc_ids = [doc_id for doc_id, _ in bm25_results]
        scores = {doc_id: score for doc_id, score in bm25_results}

        try:
            chroma_results = self.collection.get(
                ids=doc_ids, include=["documents", "metadatas"]
            )

            results = []
            for doc_id, doc, meta in zip(
                chroma_results["ids"],
                chroma_results["documents"],
                chroma_results["metadatas"],
            ):
                results.append(
                    {
                        "id": doc_id,
                        "score": scores.get(doc_id, 0.0),
                        "metadata": meta,
                        "document": doc,
                        "source": "bm25",
                    }
                )
            return results
        except Exception:
            return []

    def reciprocal_rank_fusion(
        self, results_lists: List[List[Dict]], k: int = 60
    ) -> List[Dict]:
        """
        Combine multiple ranked lists using Reciprocal Rank Fusion (RRF).

        RRF is a simple but effective fusion method that:
        - Doesn't require score normalization
        - Handles different score distributions well
        - Balances precision and recall

        Formula: RRF_score(d) = Σ 1/(k + rank_i(d))
        where k is a constant (typically 60) that dampens the impact of rank differences.
        """
        scores = defaultdict(float)
        doc_data = {}

        for results in results_lists:
            for rank, result in enumerate(results):
                doc_id = result["id"]
                scores[doc_id] += 1.0 / (k + rank + 1)

                if doc_id not in doc_data:
                    doc_data[doc_id] = result

        sorted_docs = sorted(scores.items(), key=lambda x: x[1], reverse=True)

        return [
            {**doc_data[doc_id], "rrf_score": score}
            for doc_id, score in sorted_docs
            if doc_id in doc_data
        ]

    def hybrid_retrieve(self, question: str, top_k: int = 10) -> List[Dict]:
        """
        Perform true hybrid search combining vector and BM25 retrieval.

        This addresses the limitation of the original implementation where
        "hybrid search" was just vector search with keyword filtering.
        True hybrid search runs both retrieval methods independently and
        fuses their results using RRF.
        """
        results_lists = []

        # Vector search
        vector_results = self.vector_search(question, top_k=top_k)
        if vector_results:
            results_lists.append(vector_results)

        # BM25 search
        if CONFIG.ENABLE_BM25:
            bm25_results = self.bm25_search(question, top_k=top_k)
            if bm25_results:
                results_lists.append(bm25_results)

        # Fuse results
        if len(results_lists) > 1:
            fused = self.reciprocal_rank_fusion(results_lists)
        elif results_lists:
            fused = results_lists[0]
        else:
            fused = []

        return fused[:top_k]

    def hyde_retrieve(
        self, question: str, question_lang: str, top_k: int = 10
    ) -> List[Dict]:
        """
        Perform HyDE (Hypothetical Document Embedding) retrieval.

        HyDE generates a hypothetical answer to the question, then uses that
        answer's embedding to search for similar documents. This works because
        the hypothetical answer is written in "document style" rather than
        "question style", bridging the semantic gap.

        Particularly useful for:
        - Abstract or conceptual questions
        - Questions that use different vocabulary than the documents
        - Arabic queries where question phrasing differs from document style
        """
        if not CONFIG.ENABLE_HYDE:
            return []

        # Generate hypothetical answer
        hypothetical = self.llm_client.generate_hypothetical_answer(
            question, question_lang
        )

        if not hypothetical:
            return []

        # Embed the hypothetical answer (as a passage, not query)
        hyde_embedding = self.embedder.encode(
            [f"passage: {hypothetical}"], normalize_embeddings=True
        ).tolist()

        # Search with the hypothetical embedding
        results = self.collection.query(
            query_embeddings=hyde_embedding, n_results=top_k
        )

        if not results["ids"] or not results["ids"][0]:
            return []

        return [
            {
                "id": doc_id,
                "score": 1.0 / (1 + dist),
                "metadata": meta,
                "document": doc,
                "distance": dist,
                "source": "hyde",
            }
            for doc_id, dist, meta, doc in zip(
                results["ids"][0],
                results["distances"][0],
                results["metadatas"][0],
                results["documents"][0],
            )
        ]

    def dedupe_by_parent(self, results: List[Dict]) -> List[Dict]:
        """Deduplicate results by parent_id."""
        seen_parent_ids = set()
        deduped = []
        for r in results:
            parent_id = r["metadata"].get("parent_id", r["id"])
            if parent_id in seen_parent_ids:
                continue
            seen_parent_ids.add(parent_id)
            deduped.append(r)
        return deduped

    def truncate_for_rerank(self, text: str, max_tokens: int = 256) -> str:
        """Truncate text for reranking to avoid tokenizer overflow."""
        if not text or not self.reranker:
            return text[:2000] if text else ""

        try:
            tok = self.reranker.tokenizer
            ids = tok.encode(
                text, add_special_tokens=True, truncation=True, max_length=max_tokens
            )
            return tok.decode(ids, skip_special_tokens=True)
        except Exception:
            return text[:2000]

    def rerank(
        self, query: str, candidates: List[Dict], top_n: int
    ) -> Tuple[List[Dict], List[float]]:
        """Rerank candidates using cross-encoder."""
        if not candidates or not self.reranker:
            return candidates[:top_n], []

        pool = candidates[: min(len(candidates), CONFIG.RERANK_CANDIDATES)]

        pairs = []
        for c in pool:
            doc_text = (
                c.get("parent_text")
                or c["metadata"].get("parent_text", "")
                or c.get("document", "")
            )
            doc_text = self.truncate_for_rerank(doc_text)
            pairs.append([query, doc_text])

        raw_scores = self.reranker.predict(pairs)
        raw_scores = [float(s) for s in raw_scores]

        for c, s in zip(pool, raw_scores):
            c["rerank_score"] = s
            c["rerank_prob"] = MathUtils.sigmoid(s)

        pool.sort(key=lambda x: x.get("rerank_score", -1e9), reverse=True)
        reranked = pool[: min(top_n, len(pool))]
        reranked_scores = [float(x.get("rerank_score", 0.0)) for x in reranked]

        return reranked, reranked_scores

    def apply_mmr(
        self, query: str, candidates: List[Dict], top_k: int, lambda_param: float = 0.7
    ) -> List[Dict]:
        """
        Apply Maximal Marginal Relevance (MMR) for diversity.

        MMR balances relevance to the query with diversity among selected documents.
        This prevents returning multiple chunks that say similar things.

        Formula: MMR = λ × sim(q, d) - (1-λ) × max(sim(d, d_selected))

        where:
        - λ = 1.0 means pure relevance (no diversity)
        - λ = 0.0 means pure diversity (ignore relevance)
        - λ = 0.7 (default) balances both
        """
        if not CONFIG.ENABLE_MMR or len(candidates) <= top_k:
            return candidates[:top_k]

        # Get embeddings for query and all candidates
        query_emb = self.embedder.encode(
            [f"query: {query}"], normalize_embeddings=True
        )[0]

        # Get document texts and embed them
        doc_texts = []
        for c in candidates:
            text = (
                c.get("parent_text")
                or c["metadata"].get("parent_text", "")
                or c.get("document", "")
            )
            doc_texts.append(f"passage: {text[:1000]}")  # Truncate for efficiency

        doc_embeddings = self.embedder.encode(doc_texts, normalize_embeddings=True)

        # Calculate query-document similarities
        query_sims = np.dot(doc_embeddings, query_emb)

        # MMR selection
        selected_indices = []
        remaining_indices = list(range(len(candidates)))

        while len(selected_indices) < top_k and remaining_indices:
            best_idx = None
            best_score = -float("inf")

            for idx in remaining_indices:
                # Relevance to query
                relevance = query_sims[idx]

                # Maximum similarity to already selected documents
                if selected_indices:
                    selected_embs = doc_embeddings[selected_indices]
                    candidate_emb = doc_embeddings[idx]
                    diversity_penalty = np.max(np.dot(selected_embs, candidate_emb))
                else:
                    diversity_penalty = 0

                # MMR score
                mmr_score = (
                    lambda_param * relevance - (1 - lambda_param) * diversity_penalty
                )

                if mmr_score > best_score:
                    best_score = mmr_score
                    best_idx = idx

            if best_idx is not None:
                selected_indices.append(best_idx)
                remaining_indices.remove(best_idx)

        return [candidates[i] for i in selected_indices]

    def assess_retrieval_quality(self, rerank_scores: List[float]) -> Dict[str, float]:
        """Assess retrieval quality for CRAG gating."""
        if not rerank_scores:
            return {"max_relevance": 0.0, "mean_top3": 0.0}

        probs = [MathUtils.sigmoid(s) for s in rerank_scores]
        probs_sorted = sorted(probs, reverse=True)

        max_rel = probs_sorted[0]
        top3 = probs_sorted[: min(3, len(probs_sorted))]
        mean_top3 = sum(top3) / len(top3)

        return {"max_relevance": float(max_rel), "mean_top3": float(mean_top3)}

    def should_correct(self, metrics: Dict[str, float]) -> bool:
        """Determine if CRAG corrective retrieval is needed."""
        if metrics["max_relevance"] < CONFIG.CRAG_MIN_MAX_RELEVANCE:
            return True
        if metrics["mean_top3"] < CONFIG.CRAG_MIN_MEAN_TOP3:
            return True
        return False

    def detect_comparison_query(self, question: str, question_lang: str) -> List[str]:
        """
        Detect if a question is asking to compare/contrast topics and extract individual topics.
        
        For example:
        - "compare RGB and HSV" -> ["RGB", "HSV"]
        - "what's the difference between CNN and RNN" -> ["CNN", "RNN"]
        - "قارن بين RGB و HSV" -> ["RGB", "HSV"]
        
        Returns empty list if not a comparison question.
        """
        # Comparison patterns in English
        en_patterns = [
            r"compare\s+(.+?)\s+(?:and|with|to|vs\.?|versus)\s+(.+)",
            r"difference(?:s)?\s+between\s+(.+?)\s+(?:and|&)\s+(.+)",
            r"(.+?)\s+vs\.?\s+(.+)",
            r"contrast\s+(.+?)\s+(?:and|with)\s+(.+)",
            r"(.+?)\s+compared\s+to\s+(.+)",
            r"how\s+(?:does|do|is|are)\s+(.+?)\s+differ\s+from\s+(.+)",
        ]
        
        # Comparison patterns in Arabic
        ar_patterns = [
            r"قارن\s+(?:بين\s+)?(.+?)\s+(?:و|مع)\s+(.+)",
            r"الفرق\s+بين\s+(.+?)\s+(?:و|&)\s+(.+)",
            r"ما\s+(?:هو\s+)?الفرق\s+بين\s+(.+?)\s+(?:و|&)\s+(.+)",
            r"(.+?)\s+مقابل\s+(.+)",
            r"مقارنة\s+(?:بين\s+)?(.+?)\s+(?:و|مع)\s+(.+)",
        ]
        
        patterns = ar_patterns if question_lang == "ar" else en_patterns
        
        for pattern in patterns:
            match = re.search(pattern, question, re.IGNORECASE)
            if match:
                topic1 = match.group(1).strip().strip('?.,')
                topic2 = match.group(2).strip().strip('?.,')
                # Clean up common artifacts
                topic1 = re.sub(r'^(the|a|an)\s+', '', topic1, flags=re.IGNORECASE)
                topic2 = re.sub(r'^(the|a|an)\s+', '', topic2, flags=re.IGNORECASE)
                if topic1 and topic2:
                    return [topic1, topic2]
        
        return []

    def retrieve(
        self, question: str, per_query_k: int = 5, final_k: int = None
    ) -> Tuple[List[Dict], Dict]:
        """
        Full retrieval pipeline with multi-query, HyDE, hybrid search, reranking, MMR, and CRAG.
        
        Now includes special handling for comparison queries to retrieve each topic separately.
        """
        if final_k is None:
            final_k = CONFIG.TOP_K

        debug = {
            "crag_used": False,
            "crag_metrics": None,
            "crag_metrics_after": None,
            "hyde_used": CONFIG.ENABLE_HYDE,
            "bm25_used": CONFIG.ENABLE_BM25,
            "mmr_used": CONFIG.ENABLE_MMR,
            "comparison_detected": False,
            "comparison_topics": [],
        }

        question_lang = LanguageUtils.detect_language(question)

        # Check if this is a comparison query
        comparison_topics = self.detect_comparison_query(question, question_lang)
        
        # Generate query variations
        query_variations = [question]
        query_variations.extend(
            self.llm_client.generate_multi_queries(question, question_lang)
        )
        
        # If comparison detected, add individual topic queries
        if comparison_topics:
            debug["comparison_detected"] = True
            debug["comparison_topics"] = comparison_topics
            
            # Add queries for each topic individually
            for topic in comparison_topics:
                # Add "what is X" style queries
                if question_lang == "ar":
                    query_variations.append(f"ما هو {topic}")
                    query_variations.append(f"تعريف {topic}")
                    query_variations.append(topic)
                else:
                    query_variations.append(f"what is {topic}")
                    query_variations.append(f"{topic} definition")
                    query_variations.append(f"{topic} explanation")
                    query_variations.append(topic)

        # Add translated query
        translated = self.llm_client.translate_query(question, question_lang)
        if translated:
            query_variations.append(translated)

        # Collect all results from different retrieval strategies
        all_results = []

        # 1. Hybrid retrieval for each query variation
        for q in query_variations:
            results = self.hybrid_retrieve(q, top_k=per_query_k)
            for result in results:
                result["parent_text"] = result["metadata"].get(
                    "parent_text", result["document"]
                )
                all_results.append(result)

        # 2. HyDE retrieval (only for original query to save API calls)
        if CONFIG.ENABLE_HYDE:
            hyde_results = self.hyde_retrieve(
                question, question_lang, top_k=per_query_k
            )
            for result in hyde_results:
                result["parent_text"] = result["metadata"].get(
                    "parent_text", result["document"]
                )
                all_results.append(result)

        # Sort by score and deduplicate
        all_results.sort(
            key=lambda x: x.get("rrf_score", x.get("score", 0)), reverse=True
        )
        all_results = self.dedupe_by_parent(all_results)
        candidates = all_results[: max(final_k * 2, CONFIG.RERANK_CANDIDATES)]

        # Rerank
        reranked = candidates
        rerank_scores = []

        if CONFIG.ENABLE_RERANKING and self.reranker and candidates:
            reranked, rerank_scores = self.rerank(
                question, candidates, CONFIG.RERANK_CANDIDATES
            )

        # CRAG - Corrective RAG
        if CONFIG.ENABLE_CRAG and self.reranker:
            if not rerank_scores and reranked:
                _, rerank_scores = self.rerank(
                    question, reranked, min(len(reranked), final_k)
                )

            metrics = self.assess_retrieval_quality(rerank_scores)
            debug["crag_metrics"] = metrics

            if self.should_correct(metrics):
                debug["crag_used"] = True

                # Expanded retrieval
                expanded_results = []
                for q in query_variations:
                    results = self.hybrid_retrieve(
                        q, top_k=per_query_k * CONFIG.CRAG_RETRIEVE_EXPANSION
                    )
                    for result in results:
                        result["parent_text"] = result["metadata"].get(
                            "parent_text", result["document"]
                        )
                        expanded_results.append(result)

                # Add more HyDE results
                if CONFIG.ENABLE_HYDE:
                    hyde_results = self.hyde_retrieve(
                        question,
                        question_lang,
                        top_k=per_query_k * CONFIG.CRAG_RETRIEVE_EXPANSION,
                    )
                    for result in hyde_results:
                        result["parent_text"] = result["metadata"].get(
                            "parent_text", result["document"]
                        )
                        expanded_results.append(result)

                expanded_results.sort(
                    key=lambda x: x.get("rrf_score", x.get("score", 0)), reverse=True
                )
                expanded_results = self.dedupe_by_parent(expanded_results)

                if expanded_results:
                    corrected, corrected_scores = self.rerank(
                        question, expanded_results, CONFIG.RERANK_CANDIDATES
                    )
                    corrected_metrics = self.assess_retrieval_quality(corrected_scores)
                    debug["crag_metrics_after"] = corrected_metrics

                    if (
                        corrected_metrics["max_relevance"] > metrics["max_relevance"]
                    ) or (corrected_metrics["mean_top3"] > metrics["mean_top3"]):
                        reranked = corrected

                    if (
                        self.should_correct(corrected_metrics)
                        and corrected_metrics["max_relevance"] < 0.20
                    ):
                        reranked = []

        # Apply MMR for diversity
        if CONFIG.ENABLE_MMR and reranked:
            reranked = self.apply_mmr(
                question,
                reranked,
                top_k=CONFIG.MMR_TOP_K,
                lambda_param=CONFIG.MMR_LAMBDA,
            )
        else:
            reranked = reranked[:final_k]

        return reranked, debug


# ============================================================================
# SEMANTIC CACHE
# ============================================================================
class SemanticCache:
    """Semantic caching for query results."""

    def __init__(self, cache: OrderedDict, embedder: SentenceTransformer):
        self.cache = cache
        self.embedder = embedder

    @staticmethod
    def make_scope_key(
        provider_key: str, model_name: str, collection_count: int
    ) -> str:
        """Create a scope key for cache isolation."""
        return f"{provider_key}::{model_name}::count:{collection_count}"

    def lookup(
        self, query: str, query_embedding: List[float], scope_key: str
    ) -> Optional[Dict]:
        """Look up a semantically similar cached entry."""
        if not CONFIG.ENABLE_SEMANTIC_CACHE:
            return None

        best_key = None
        best_sim = -1.0
        best_payload = None

        for k, payload in self.cache.items():
            if not k.startswith(scope_key + "::"):
                continue

            emb = payload.get("query_embedding")
            if not emb:
                continue

            sim = MathUtils.cosine_sim(query_embedding, emb)
            if sim > best_sim:
                best_sim = sim
                best_key = k
                best_payload = payload

        if best_payload and best_sim >= CONFIG.CACHE_SIM_THRESHOLD:
            self.cache.move_to_end(best_key, last=True)
            result = dict(best_payload)
            result["cache_similarity"] = float(best_sim)
            return result

        return None

    def store(
        self,
        query: str,
        query_embedding: List[float],
        answer: str,
        sources: List[Dict],
        followups: List[str],
        question_lang: str,
        scope_key: str,
    ):
        """Store a query result in cache."""
        if not CONFIG.ENABLE_SEMANTIC_CACHE:
            return

        qh = hashlib.md5(query.encode("utf-8")).hexdigest()[:12]
        k = f"{scope_key}::{qh}"

        self.cache[k] = {
            "query": query,
            "query_embedding": query_embedding,
            "answer": answer,
            "sources": sources,
            "followups": followups,
            "question_lang": question_lang,
            "cached_at": time.time(),
        }
        self.cache.move_to_end(k, last=True)

        # LRU eviction
        while len(self.cache) > CONFIG.CACHE_MAX_ENTRIES:
            self.cache.popitem(last=False)


# ============================================================================
# ANSWER GENERATION - FIXED VERSION
# ============================================================================
class AnswerGenerator:
    """Generate answers with proper citations, context compression, and language handling."""

    def __init__(self, llm_client: LLMClient):
        self.llm_client = llm_client

    def validate_arabic_response(self, text: str) -> str:
        """Validate and clean Arabic responses."""
        # Remove Chinese characters
        cleaned = re.sub(r"[\u4e00-\u9fff]", "", text)

        # Remove Vietnamese diacritics
        vietnamese_pattern = (
            r"[àáảãạăằắẳẵặâầấẩẫậèéẻẽẹêềếểễệìíỉĩịòóỏõọôồốổỗộơờớởỡợùúủũụưừứửữựỳýỷỹỵđ]"
        )
        cleaned = re.sub(vietnamese_pattern, "", cleaned, flags=re.IGNORECASE)

        # Check for excessive Latin text outside parentheses
        text_outside_parens = re.sub(r"\([^)]*\)", "", cleaned)
        latin_outside = re.findall(
            r"(?<!\()[a-zA-Z]{3,}(?![^(]*\))", text_outside_parens
        )

        if latin_outside and len(latin_outside) > 3:
            unique_terms = list(set(latin_outside))[:5]
            st.warning(
                f"⚠️ Language mixing detected. Terms found: {', '.join(unique_terms)}"
            )

        return cleaned

    def compress_contexts(
        self, retrieved_docs: List[Dict], question: str, question_lang: str
    ) -> List[Dict]:
        """
        Compress retrieved contexts by extracting only relevant sentences.

        This reduces noise in the context sent to the generator, improving
        answer quality and reducing hallucination risk.
        """
        if not CONFIG.ENABLE_CONTEXT_COMPRESSION:
            return retrieved_docs

        compressed_docs = []
        for doc in retrieved_docs:
            parent_text = doc.get("parent_text", doc.get("document", ""))

            # Compress the context
            compressed = self.llm_client.compress_context(
                question, parent_text, question_lang, CONFIG.COMPRESSION_MAX_SENTENCES
            )

            # Create new doc with compressed text
            compressed_doc = dict(doc)
            compressed_doc["original_parent_text"] = parent_text
            compressed_doc["parent_text"] = compressed
            compressed_docs.append(compressed_doc)

        return compressed_docs

    def get_arabic_system_prompt(self) -> str:
        """Get the system prompt for Arabic responses - FIXED VERSION."""
        return """أنت مدرس جامعي متخصص في مجالات هندسة الذكاء الاصطناعي (الرؤية الحاسوبية، معالجة اللغات الطبيعية، التعلم الآلي).

═══════════════════════════════════════════════════════════════
قاعدة حاسمة - المصادر المقدمة فقط
═══════════════════════════════════════════════════════════════
يجب أن تستخدم فقط المعلومات من المصادر المقدمة. لا تستخدم معرفتك العامة أبداً.
- إذا لم تكن المعلومات في المصادر: قل "هذه المعلومات غير متوفرة في مواد المحاضرات المقدمة."
- إذا كانت المعلومات متوفرة جزئياً: أجب فقط بما تحتويه المصادر، استشهد به، واذكر ما هو ناقص.
- لا تقل أبداً "بناءً على المعرفة العامة" أو تجب بدون استشهادات.

═══════════════════════════════════════════════════════════════
التعامل مع أسئلة المقارنة/التوليف
═══════════════════════════════════════════════════════════════
إذا طُلب منك مقارنة أو ربط موضوعات (مثل: "قارن بين RGB و HSV"):
1. تحقق مما إذا كان كل موضوع موصوفاً في المصادر (ربما في مصادر مختلفة)
2. إذا كان الموضوع أ في [1] والموضوع ب في [3]: قم بالتوليف بشرح كل منهما من مصدره مع الاستشهادات الصحيحة
3. إذا وُجد موضوع واحد فقط: اشرح ما وجدته واذكر أن الآخر غير موجود في المواد
4. إذا لم يُوجد أي منهما: اذكر بوضوح أن هذه المعلومات غير موجودة في المواد المقدمة

═══════════════════════════════════════════════════════════════
قواعد الاستشهاد - إلزامية
═══════════════════════════════════════════════════════════════
- استخدم استشهادات مثل [1]، [2]، ... التي تشير إلى المصادر المقدمة.
- كل ادعاء واقعي يجب أن يكون له استشهاد. بدون استثناءات.
- لا تخترع مصادر أو معلومات.

═══════════════════════════════════════════════════════════════
قواعد اللغة - إلزامية
═══════════════════════════════════════════════════════════════
- اكتب بالعربية الفصحى فقط
- المصطلحات التقنية: اكتب بالعربية أولاً ثم (English) بين قوسين
- مثال صحيح: "الرؤية الحاسوبية (Computer Vision)"
- مثال خاطئ: "Computer Vision هو مجال" أو "الـ processing"

═══════════════════════════════════════════════════════════════
تنسيق الإجابة
═══════════════════════════════════════════════════════════════
- اكتب فقرات متصلة (ليس نقاط)
- استخدم علامات الترقيم العربية: ، ؛ ؟
- أسلوب أكاديمي واضح
- 150-250 كلمة تقريباً (أقصر إذا كانت المواد المصدرية محدودة)"""

    def get_english_system_prompt(self) -> str:
        """Get the system prompt for English responses - FIXED VERSION."""
        return """You are a university tutor specializing in AI Engineering (Computer Vision, NLP, Machine Learning).

═══════════════════════════════════════════════════════════════
CRITICAL RULE - SOURCE MATERIAL ONLY
═══════════════════════════════════════════════════════════════
You must ONLY use information from the provided sources. NEVER use your general knowledge.
- If information is NOT in the sources: Say "This information is not available in the provided lecture materials."
- If information IS partially available: Answer ONLY what the sources contain, cite it, and note what's missing.
- NEVER say things like "based on general knowledge" or answer without citations.

═══════════════════════════════════════════════════════════════
HANDLING COMPARISON/SYNTHESIS QUESTIONS
═══════════════════════════════════════════════════════════════
If asked to compare, contrast, or relate topics (e.g., "compare RGB and HSV"):
1. Check if EACH topic is described in the sources (possibly in different sources)
2. If Topic A is in [1] and Topic B is in [3]: Synthesize by explaining each from its source with proper citations
3. If only ONE topic is found: Explain what you found and state the other is not in the materials
4. If NEITHER is found: State clearly that this information is not in the provided materials

═══════════════════════════════════════════════════════════════
CITATION REQUIREMENTS - MANDATORY
═══════════════════════════════════════════════════════════════
- Use citations like [1], [2], ... that refer to the provided sources.
- EVERY factual claim MUST have a citation. No exceptions.
- Do not invent sources or information.

═══════════════════════════════════════════════════════════════
RESPONSE FORMAT
═══════════════════════════════════════════════════════════════
- Write in clear, flowing paragraphs (NOT bullet points)
- Academic but accessible language
- Translate any Arabic content from sources to English
- Aim for 150-250 words (shorter if limited source material)"""

    def generate_answer(
        self, retrieved_docs: List[Dict], question: str
    ) -> Tuple[str, List[Dict]]:
        """Generate an answer with citations and optional context compression."""
        question_lang = LanguageUtils.detect_language(question)

        if not retrieved_docs:
            no_info = (
                "لا تتوفر معلومات ذات صلة في المستندات المرفوعة. الرجاء التأكد من رفع مواد المحاضرات ذات الصلة."
                if question_lang == "ar"
                else "I couldn't find relevant information in the uploaded documents. Please make sure to upload the relevant lecture materials."
            )
            return no_info, []

        # Apply context compression if enabled
        docs_to_use = self.compress_contexts(retrieved_docs, question, question_lang)

        # Build context from parent texts
        context_parts = []
        sources = []

        for i, doc in enumerate(docs_to_use, 1):
            parent_text = doc.get("parent_text", doc.get("document", ""))
            context_parts.append(f"[Source {i}]:\n{parent_text}")
            meta = doc["metadata"]
            sources.append(
                {
                    "id": doc["id"],                 # << key for evaluation
                    "source": meta.get("source", ""),
                    "page": meta.get("page", ""),
                    "type": meta.get("type", ""),
                    "lang": meta.get("lang", "en"),
                    "parent_id": meta.get("parent_id", ""),
                    "used_text": parent_text,        # << used for faithfulness eval
                }
            )


        context = "\n\n".join(context_parts)

        # Select system prompt based on language
        if question_lang == "ar":
            system_prompt = self.get_arabic_system_prompt()
            reminder = """تذكير مهم: 
1. اكتب بالعربية فقط، مع المصطلحات التقنية بين قوسين بالإنجليزية
2. استشهد بكل معلومة من المصادر
3. إذا لم تجد المعلومة في المصادر، قل ذلك بوضوح
4. لا تستخدم معرفتك العامة أبداً"""
        else:
            system_prompt = self.get_english_system_prompt()
            reminder = """CRITICAL REMINDER: 
1. ONLY use information from the sources above
2. Cite every fact with [1], [2], etc.
3. If the information is not in the sources, say so clearly
4. NEVER use your general knowledge - if it's not cited, don't say it"""

        messages = [
            {"role": "system", "content": system_prompt},
            {
                "role": "user",
                "content": f"""LECTURE MATERIAL SOURCES:
{context}

STUDENT QUESTION: {question}

{reminder}""",
            },
        ]

        # Lower temperature for more factual responses
        temperature = 0.2 if question_lang == "ar" else 0.3

        answer = self.llm_client.chat(messages, temperature=temperature)

        # Validate Arabic responses
        if question_lang == "ar":
            answer = self.validate_arabic_response(answer)

            # Check citation count
            citations = re.findall(r"\[\d+\]", answer)
            if len(citations) < 1 and "غير متوفرة" not in answer and "غير موجودة" not in answer:
                st.warning(
                    "⚠️ Response has no citations - the information may not be in the source materials"
                )

        # Check English responses for missing citations
        if question_lang == "en":
            citations = re.findall(r"\[\d+\]", answer)
            if len(citations) < 1 and "not available" not in answer.lower() and "not in the" not in answer.lower():
                st.warning(
                    "⚠️ Response has no citations - the information may not be in the source materials"
                )

        if question_lang == "ar":
            header = "المصادر:"
            lines = [
                f"[{i}] {s.get('source','')} - صفحة {s.get('page','')}"
                for i, s in enumerate(sources, 1)
            ]
        else:
            header = "Sources:"
            lines = [
                f"[{i}] {s.get('source','')} - page {s.get('page','')}"
                for i, s in enumerate(sources, 1)
            ]

        answer = answer.strip() + "\n\n" + header + "\n" + "\n".join(lines)


        return answer, sources

    def generate_followups(
        self, question: str, answer: str, question_lang: str
    ) -> List[str]:
        """Generate follow-up questions."""
        if question_lang == "ar":
            instruction = """بناءً على هذا السؤال والجواب، أنشئ 3 أسئلة متابعة قصيرة بالعربية.
ركز على: توضيح المفاهيم، مواضيع ذات صلة، أمثلة تطبيقية.
اكتب الأسئلة فقط، سؤال في كل سطر، بدون ترقيم."""
        else:
            instruction = """Based on this Q&A, generate 3 short follow-up questions in English.
Focus on: clarifying concepts, related topics, practical examples.
Output ONLY the questions, one per line, without numbering."""

        try:
            result = self.llm_client.chat(
                [
                    {
                        "role": "user",
                        "content": f"{instruction}\n\nQuestion: {question}\n\nAnswer: {answer[:500]}...",
                    }
                ],
                temperature=0.8,
            )

            lines = [l.strip() for l in result.strip().split("\n") if l.strip()]
            cleaned = []
            for line in lines[:3]:
                line = re.sub(r"^[\d\.\-\*\)]+\s*", "", line)
                if line and len(line) > 10:
                    cleaned.append(line)
            return cleaned
        except Exception:
            return []


# ============================================================================
# EVALUATION FRAMEWORK
# ============================================================================
@dataclass
class EvalQuery:
    """A single evaluation query with ground truth."""

    query: str
    language: str  # "ar" or "en"
    ground_truth_answer: str
    relevant_doc_ids: List[str]  # IDs of documents that should be retrieved
    category: str = "general"  # e.g., "factual", "conceptual", "multi-hop"


@dataclass
class EvalResult:
    """Result of evaluating a single query."""

    query: str
    language: str
    retrieved_doc_ids: List[str]
    answer: str

    # Retrieval metrics
    hit_rate: float  # 1 if any relevant doc retrieved, 0 otherwise
    recall_at_k: float  # fraction of relevant docs retrieved
    precision_at_k: float  # fraction of retrieved docs that are relevant
    mrr: float  # Mean Reciprocal Rank

    # Generation metrics (require LLM-as-judge)
    faithfulness: Optional[float] = None
    answer_relevancy: Optional[float] = None

    # Metadata
    latency_ms: float = 0.0
    cache_hit: bool = False


class RAGEvaluator:
    """
    Evaluation framework for the RAG pipeline.

    Implements RAGAS-style metrics:
    - Retrieval: Hit Rate, Recall@K, Precision@K, MRR
    - Generation: Faithfulness, Answer Relevancy (using LLM-as-judge)
    """

    def __init__(self, rag_pipeline, llm_client: LLMClient):
        self.rag_pipeline = rag_pipeline
        self.llm_client = llm_client

    def evaluate_retrieval(
        self, retrieved_ids: List[str], relevant_ids: List[str]
    ) -> Dict[str, float]:
        """
        Calculate retrieval metrics.

        Returns dict with: hit_rate, recall_at_k, precision_at_k, mrr
        """
        if not relevant_ids:
            return {
                "hit_rate": 0.0,
                "recall_at_k": 0.0,
                "precision_at_k": 0.0,
                "mrr": 0.0,
            }

        relevant_set = set(relevant_ids)

        # Hit rate: 1 if any relevant doc in retrieved
        hits = [1 if doc_id in relevant_set else 0 for doc_id in retrieved_ids]
        hit_rate = 1.0 if any(hits) else 0.0

        # Recall@K: fraction of relevant docs that were retrieved
        retrieved_relevant = sum(
            1 for doc_id in retrieved_ids if doc_id in relevant_set
        )
        recall_at_k = retrieved_relevant / len(relevant_ids)

        # Precision@K: fraction of retrieved that are relevant
        precision_at_k = (
            retrieved_relevant / len(retrieved_ids) if retrieved_ids else 0.0
        )

        # MRR: 1 / position of first relevant doc
        mrr = 0.0
        for i, doc_id in enumerate(retrieved_ids, 1):
            if doc_id in relevant_set:
                mrr = 1.0 / i
                break

        return {
            "hit_rate": hit_rate,
            "recall_at_k": recall_at_k,
            "precision_at_k": precision_at_k,
            "mrr": mrr,
        }

    def evaluate_faithfulness(
        self, question: str, answer: str, context: str, lang: str
    ) -> float:
        """
        Evaluate if the answer is faithful to (grounded in) the context.

        Uses LLM-as-judge to assess whether claims in the answer are supported
        by the provided context.
        """
        if lang == "ar":
            prompt = f"""قيّم مدى التزام الإجابة بالمصادر المقدمة.

السؤال: {question}

المصادر:
{context[:2000]}

الإجابة: {answer}

أعط درجة من 0 إلى 1:
- 1.0: كل المعلومات في الإجابة موجودة في المصادر
- 0.5: بعض المعلومات مدعومة وبعضها لا
- 0.0: الإجابة لا تتوافق مع المصادر

اكتب الرقم فقط (مثال: 0.8)"""
        else:
            prompt = f"""Evaluate if the answer is faithful to (grounded in) the provided context.

Question: {question}

Context:
{context[:2000]}

Answer: {answer}

Give a score from 0 to 1:
- 1.0: All information in the answer is supported by the context
- 0.5: Some information is supported, some is not
- 0.0: The answer contradicts or is not grounded in the context

Output ONLY the number (e.g., 0.8)"""

        try:
            result = self.llm_client.chat(
                [{"role": "user", "content": prompt}],
                temperature=0.1,
                max_tokens=10,
            )
            score = float(re.search(r"[0-9.]+", result).group())
            return min(1.0, max(0.0, score))
        except Exception:
            return 0.5  # Default to neutral if parsing fails

    def evaluate_answer_relevancy(self, question: str, answer: str, lang: str) -> float:
        """
        Evaluate if the answer actually addresses the question.

        Uses LLM-as-judge to assess relevancy.
        """
        if lang == "ar":
            prompt = f"""قيّم مدى إجابة الرد على السؤال المطروح.

السؤال: {question}

الإجابة: {answer}

أعط درجة من 0 إلى 1:
- 1.0: الإجابة تجيب مباشرة وبشكل كامل على السؤال
- 0.5: الإجابة جزئية أو غير مباشرة
- 0.0: الإجابة لا تتعلق بالسؤال

اكتب الرقم فقط (مثال: 0.8)"""
        else:
            prompt = f"""Evaluate if the answer addresses the question asked.

Question: {question}

Answer: {answer}

Give a score from 0 to 1:
- 1.0: The answer directly and completely addresses the question
- 0.5: The answer is partial or indirect
- 0.0: The answer does not relate to the question

Output ONLY the number (e.g., 0.8)"""

        try:
            result = self.llm_client.chat(
                [{"role": "user", "content": prompt}],
                temperature=0.1,
                max_tokens=10,
            )
            score = float(re.search(r"[0-9.]+", result).group())
            return min(1.0, max(0.0, score))
        except Exception:
            return 0.5

    def evaluate_single(self, eval_query: EvalQuery) -> EvalResult:
        """Evaluate a single query."""
        start_time = time.time()

        # Run the RAG pipeline
        answer, sources, followups, debug = self.rag_pipeline.answer(eval_query.query)

        latency_ms = (time.time() - start_time) * 1000

        # Get retrieved doc IDs
        retrieved_ids = [s.get("id", "") for s in sources if s.get("id")]

        # Calculate retrieval metrics
        retrieval_metrics = self.evaluate_retrieval(
            retrieved_ids, eval_query.relevant_doc_ids
        )

        # Build context string for faithfulness evaluation
        context_parts = []
        for i, s in enumerate(sources, 1):
            context_parts.append(f"[{i}] {s.get('used_text', '')[:800]}")
        context = "\n\n".join(context_parts)

        # Calculate generation metrics
        faithfulness = self.evaluate_faithfulness(
            eval_query.query, answer, context, eval_query.language
        )
        answer_relevancy = self.evaluate_answer_relevancy(
            eval_query.query, answer, eval_query.language
        )

        return EvalResult(
            query=eval_query.query,
            language=eval_query.language,
            retrieved_doc_ids=retrieved_ids,
            answer=answer,
            hit_rate=retrieval_metrics["hit_rate"],
            recall_at_k=retrieval_metrics["recall_at_k"],
            precision_at_k=retrieval_metrics["precision_at_k"],
            mrr=retrieval_metrics["mrr"],
            faithfulness=faithfulness,
            answer_relevancy=answer_relevancy,
            latency_ms=latency_ms,
            cache_hit=debug.get("cache_hit", False),
        )

    def evaluate_dataset(self, queries: List[EvalQuery]) -> Dict:
        """
        Evaluate a full dataset of queries.

        Returns aggregate metrics and per-query results.
        """
        results = []

        for query in queries:
            result = self.evaluate_single(query)
            results.append(result)

        # Aggregate metrics
        n = len(results)
        if n == 0:
            return {"error": "No queries to evaluate"}

        # Split by language for language-specific metrics
        ar_results = [r for r in results if r.language == "ar"]
        en_results = [r for r in results if r.language == "en"]

        def aggregate(result_list):
            if not result_list:
                return {}
            n = len(result_list)
            return {
                "hit_rate": sum(r.hit_rate for r in result_list) / n,
                "recall_at_k": sum(r.recall_at_k for r in result_list) / n,
                "precision_at_k": sum(r.precision_at_k for r in result_list) / n,
                "mrr": sum(r.mrr for r in result_list) / n,
                "faithfulness": sum(r.faithfulness or 0 for r in result_list) / n,
                "answer_relevancy": sum(r.answer_relevancy or 0 for r in result_list)
                / n,
                "avg_latency_ms": sum(r.latency_ms for r in result_list) / n,
                "cache_hit_rate": sum(1 for r in result_list if r.cache_hit) / n,
            }

        return {
            "overall": aggregate(results),
            "arabic": aggregate(ar_results),
            "english": aggregate(en_results),
            "num_queries": n,
            "num_arabic": len(ar_results),
            "num_english": len(en_results),
            "timestamp": datetime.now().isoformat(),
            "results": [
                {
                    "query": r.query,
                    "language": r.language,
                    "hit_rate": r.hit_rate,
                    "recall_at_k": r.recall_at_k,
                    "faithfulness": r.faithfulness,
                    "answer_relevancy": r.answer_relevancy,
                    "latency_ms": r.latency_ms,
                }
                for r in results
            ],
        }

def create_eval_dataset_from_corpus(collection, llm_client: LLMClient, n: int = 5) -> List[EvalQuery]:
        """
        Auto-generate evaluation queries from the currently ingested corpus.
        Each query has a known relevant chunk id, so retrieval metrics are meaningful.
        """
        try:
            count = collection.count()
        except Exception:
            count = 0

        if count <= 0:
            return []

        # Get IDs (ok for homework-scale corpora)
        all_ids = collection.get(include=[])["ids"]
        sample_ids = random.sample(all_ids, k=min(n, len(all_ids)))

        payload = collection.get(ids=sample_ids, include=["documents", "metadatas"])

        eval_queries: List[EvalQuery] = []
        for chunk_id, doc, meta in zip(payload["ids"], payload["documents"], payload["metadatas"]):
            lang = meta.get("lang", "en")
            excerpt = (meta.get("parent_text") or doc or "").strip()
            excerpt = excerpt[:1500]

            if not excerpt:
                continue

            if lang == "ar":
                prompt = f"""
    أنت تُنشئ أسئلة تقييم لنظام RAG.
    اقرأ المقتطف، وأنشئ سؤالًا واحدًا يمكن الإجابة عنه من هذا المقتطف فقط، ثم إجابة مرجعية قصيرة.
    أخرج JSON فقط بالشكل:
    {{"question":"...","answer":"..."}}

    المقتطف:
    {excerpt}
    """
            else:
                prompt = f"""
    You are creating evaluation questions for a RAG system.
    Read the excerpt and write ONE question that can be answered using ONLY this excerpt, plus a short reference answer.
    Output STRICT JSON only:
    {{"question":"...","answer":"..."}}

    Excerpt:
    {excerpt}
    """

            raw = llm_client.chat([{"role": "user", "content": prompt}], temperature=0.3, max_tokens=220)

            # Parse JSON robustly
            m = re.search(r"\{.*\}", raw, flags=re.DOTALL)
            if not m:
                continue
            try:
                obj = json.loads(m.group(0))
                q = (obj.get("question") or "").strip()
                a = (obj.get("answer") or "").strip()
                if len(q) < 8 or len(a) < 5:
                    continue
            except Exception:
                continue

            eval_queries.append(
                EvalQuery(
                    query=q,
                    language=lang,
                    ground_truth_answer=a,
                    relevant_doc_ids=[chunk_id],   # this is the key fix
                    category="auto",
                )
            )

        return eval_queries

def create_sample_eval_dataset() -> List[EvalQuery]:
    """
    Create a sample evaluation dataset.

    In practice, you would want 100+ queries covering:
    - Different languages (Arabic/English)
    - Different question types (factual, conceptual, multi-hop)
    - Different difficulty levels
    - Edge cases
    """
    return [
        # English queries
        EvalQuery(
            query="What is RAG and why is it used?",
            language="en",
            ground_truth_answer="RAG (Retrieval-Augmented Generation) combines LLMs with external knowledge retrieval to ground responses in evidence, reducing hallucination.",
            relevant_doc_ids=[],  # Fill with actual doc IDs from your corpus
            category="conceptual",
        ),
        EvalQuery(
            query="What are the different chunking strategies?",
            language="en",
            ground_truth_answer="Chunking strategies include fixed size, recursive, semantic, sentence-based, document-level, and agentic chunking.",
            relevant_doc_ids=[],
            category="factual",
        ),
        EvalQuery(
            query="How does hybrid search work?",
            language="en",
            ground_truth_answer="Hybrid search combines vector similarity search with keyword-based BM25 search, fusing results using methods like RRF.",
            relevant_doc_ids=[],
            category="conceptual",
        ),
        # Arabic queries
        EvalQuery(
            query="ما هو RAG ولماذا يستخدم؟",
            language="ar",
            ground_truth_answer="RAG هو توليد معزز بالاسترجاع يجمع بين نماذج اللغة الكبيرة واسترجاع المعرفة الخارجية لتأسيس الردود على أدلة.",
            relevant_doc_ids=[],
            category="conceptual",
        ),
        EvalQuery(
            query="ما هي استراتيجيات تقسيم النص؟",
            language="ar",
            ground_truth_answer="تشمل استراتيجيات التقسيم: الحجم الثابت، التقسيم المتكرر، الدلالي، على مستوى الجملة، على مستوى المستند.",
            relevant_doc_ids=[],
            category="factual",
        ),
    ]


# ============================================================================
# RAG PIPELINE
# ============================================================================
class RAGPipeline:
    """Main RAG pipeline orchestrating all components."""

    def __init__(
        self, llm_client: LLMClient, embedder: SentenceTransformer, collection
    ):
        self.llm_client = llm_client
        self.embedder = embedder
        self.collection = collection
        self.retriever = Retriever(embedder, collection, llm_client)
        self.answer_generator = AnswerGenerator(llm_client)
        self.cache = SemanticCache(st.session_state.semantic_cache, embedder)

    def answer(self, question: str) -> Tuple[str, List[Dict], List[str], Dict]:
        """Process a question and return answer, sources, followups, and debug info."""
        question_lang = LanguageUtils.detect_language(question)
        debug = {
            "cache_hit": False,
            "cache_similarity": None,
            "crag_used": False,
            "crag_metrics": None,
            "hyde_used": CONFIG.ENABLE_HYDE,
            "bm25_used": CONFIG.ENABLE_BM25,
            "mmr_used": CONFIG.ENABLE_MMR,
            "compression_used": CONFIG.ENABLE_CONTEXT_COMPRESSION,
        }

        # Get query embedding
        query_emb = self.embedder.encode(
            [f"query: {question}"], normalize_embeddings=True
        ).tolist()[0]

        # Create cache scope key
        try:
            collection_count = self.collection.count()
        except Exception:
            collection_count = 0

        scope_key = SemanticCache.make_scope_key(
            st.session_state.selected_provider, self.llm_client.model, collection_count
        )

        # Check cache
        if CONFIG.ENABLE_SEMANTIC_CACHE:
            hit = self.cache.lookup(question, query_emb, scope_key)
            if hit:
                debug["cache_hit"] = True
                debug["cache_similarity"] = hit.get("cache_similarity")
                return (
                    hit["answer"],
                    hit.get("sources", []),
                    hit.get("followups", []),
                    debug,
                )

        # Retrieve documents
        retrieved_docs, retrieval_debug = self.retriever.retrieve(question)
        debug.update(retrieval_debug)

        # Generate answer
        answer, sources = self.answer_generator.generate_answer(
            retrieved_docs, question
        )

        # Generate follow-ups
        followups = self.answer_generator.generate_followups(
            question, answer, question_lang
        )

        # Store in cache
        if CONFIG.ENABLE_SEMANTIC_CACHE:
            self.cache.store(
                query=question,
                query_embedding=query_emb,
                answer=answer,
                sources=sources,
                followups=followups,
                question_lang=question_lang,
                scope_key=scope_key,
            )

        return answer, sources, followups, debug


# ============================================================================
# UI COMPONENTS
# ============================================================================
class UIRenderer:
    """Render UI components."""

    @staticmethod
    def render_source_cards(metadatas: List[Dict]):
        """Render source cards."""
        if not metadatas:
            return

        st.markdown(
            '<div class="section-label">📚 Sources</div>', unsafe_allow_html=True
        )

        cols = st.columns(min(len(metadatas), 3))

        for i, meta in enumerate(metadatas):
            col_idx = i % 3
            with cols[col_idx]:
                source_name = meta.get("source", "Unknown")
                page = meta.get("page", "?")
                lang = meta.get("lang", "en")

                display_name = (
                    source_name[:25] + "..." if len(source_name) > 25 else source_name
                )
                page_label = "صفحة" if lang == "ar" else "Page"

                st.markdown(
                    f"""
                <div class="source-card">
                    <div class="source-number">{i + 1}</div>
                    <div style="overflow: hidden;">
                        <div style="font-weight: 500; font-size: 0.85rem; color: #1a1a1a;
                                    white-space: nowrap; overflow: hidden; text-overflow: ellipsis;">
                            {display_name}
                        </div>
                        <div style="color: #666; font-size: 0.75rem;">{page_label} {page}</div>
                    </div>
                </div>
                """,
                    unsafe_allow_html=True,
                )

    @staticmethod
    def render_answer(answer: str, is_rtl: bool = False):
        """Render the answer with styled citations."""

        def replace_citation(match):
            return f'<span class="citation">{match.group(1)}</span>'

        styled = re.sub(r"\[(\d+)\]", replace_citation, answer)
        styled = styled.replace("\n\n", "</p><p>").replace("\n", "<br>")
        styled = f"<p>{styled}</p>"

        css_class = "assistant-message-rtl" if is_rtl else "assistant-message"
        st.markdown(f'<div class="{css_class}">{styled}</div>', unsafe_allow_html=True)

    @staticmethod
    def render_followups(
        questions: List[str], question_lang: str, message_idx: int = 0
    ) -> Optional[str]:
        """Render follow-up question buttons."""
        if not questions:
            return None

        label = "أسئلة ذات صلة" if question_lang == "ar" else "Related Questions"
        st.markdown(
            f'<div class="section-label">💡 {label}</div>', unsafe_allow_html=True
        )

        cols = st.columns(len(questions))
        for i, (col, q) in enumerate(zip(cols, questions)):
            with col:
                # Use message_idx + question index + timestamp-based unique ID to ensure uniqueness
                unique_key = f"followup_msg{message_idx}_q{i}_{hash(q)}_{id(q)}"
                if st.button(q, key=unique_key, use_container_width=True):
                    return q
        return None

    @staticmethod
    def render_eval_metrics(metrics: Dict):
        """Render evaluation metrics in a nice format."""
        st.markdown("### 📊 Evaluation Results")

        col1, col2, col3 = st.columns(3)

        with col1:
            st.markdown("**Retrieval Metrics**")
            st.metric("Hit Rate", f"{metrics.get('hit_rate', 0):.2%}")
            st.metric("Recall@K", f"{metrics.get('recall_at_k', 0):.2%}")
            st.metric("MRR", f"{metrics.get('mrr', 0):.3f}")

        with col2:
            st.markdown("**Generation Metrics**")
            st.metric("Faithfulness", f"{metrics.get('faithfulness', 0):.2%}")
            st.metric("Answer Relevancy", f"{metrics.get('answer_relevancy', 0):.2%}")

        with col3:
            st.markdown("**Performance**")
            st.metric("Avg Latency", f"{metrics.get('avg_latency_ms', 0):.0f} ms")
            st.metric("Cache Hit Rate", f"{metrics.get('cache_hit_rate', 0):.2%}")

    @classmethod
    def render_message(
        cls,
        role: str,
        content: str,
        sources: List[Dict] = None,
        followups: List[str] = None,
        question_lang: str = "en",
        message_idx: int = 0,
    ) -> Optional[str]:
        """Render a chat message."""
        if role == "user":
            is_rtl = LanguageUtils.detect_language(content) == "ar"
            direction = 'dir="rtl"' if is_rtl else ""
            st.markdown(
                f'<div class="user-message" {direction}>{content}</div>',
                unsafe_allow_html=True,
            )
        else:
            if sources:
                cls.render_source_cards(sources)
            cls.render_answer(content, LanguageUtils.detect_language(content) == "ar")
            if followups:
                return cls.render_followups(followups, question_lang, message_idx)
        return None


# ============================================================================
# SIDEBAR
# ============================================================================
def render_sidebar():
    """Render the sidebar with settings and document upload."""
    with st.sidebar:
        st.markdown("### ⚙️ LLM Provider")

        selected = st.selectbox(
            "Choose provider",
            options=list(LLM_PROVIDERS.keys()),
            index=list(LLM_PROVIDERS.keys()).index(st.session_state.selected_provider),
            label_visibility="collapsed",
        )
        st.session_state.selected_provider = selected

        provider_config = LLM_PROVIDERS[selected]
        st.caption(f"ℹ️ {provider_config['notes']}")

        api_key = st.text_input(
            "API Key",
            type="password",
            value=st.session_state.api_key,
            placeholder=f"Enter {provider_config['name']} API key",
        )
        if api_key:
            st.session_state.api_key = api_key

        st.markdown(f"[🔑 Get free API key]({provider_config['get_key_url']})")

        # Test connection button
        if st.session_state.api_key:
            if st.button("🔌 Test Connection", use_container_width=True):
                with st.spinner("Testing connection..."):
                    test_client = LLMClient(
                        st.session_state.selected_provider, st.session_state.api_key
                    )
                    success, message = test_client.test_connection()
                    if success:
                        st.success(f"✅ {message}")
                    else:
                        st.error(f"❌ {message}")

        st.markdown("---")
        st.markdown("### 📄 Documents")

        uploaded_files = st.file_uploader(
            "Upload files",
            type=["pdf", "ppt", "pptx"],
            accept_multiple_files=True,
            label_visibility="collapsed",
        )

        if st.button("📥 Ingest", use_container_width=True):
            if not uploaded_files:
                st.error("Upload files first")
            else:
                embedder = load_embedder()
                _, collection = get_chroma_collection()
                ingester = DocumentIngester(embedder, collection)

                total = 0
                for file in uploaded_files:
                    with st.spinner(f"Processing {file.name}..."):
                        added, status = ingester.ingest(file)
                        if status == "exists":
                            st.info(f"✓ {file.name} already indexed")
                        elif status == "success":
                            st.success(f"✓ {file.name}: {added} chunks added")
                            total += added
                        else:
                            st.error(f"✗ {file.name}: {status}")

                if total > 0:
                    st.success(f"✅ Total: {total} new chunks indexed")

        # Show document count
        try:
            _, collection = get_chroma_collection()
            count = collection.count()
            st.caption(f"📊 {count} chunks in database")
        except Exception:
            pass

        st.markdown("---")
        st.markdown("### 🔧 Advanced Settings")

        with st.expander("Retrieval Settings"):
            st.caption("**Chunking:**")
            st.caption(f"• Parent: {CONFIG.PARENT_CHUNK_SIZE} tokens")
            st.caption(f"• Child: {CONFIG.CHILD_CHUNK_SIZE} tokens")
            st.caption(f"• Top-K: {CONFIG.TOP_K}")

            st.caption("**Search:**")
            st.caption(f"• Vector Search: ✅ Enabled")
            st.caption(f"• BM25 Search: {'✅' if CONFIG.ENABLE_BM25 else '❌'}")
            st.caption(f"• HyDE: {'✅' if CONFIG.ENABLE_HYDE else '❌'}")
            st.caption(f"• Multi-Query: {CONFIG.MULTI_QUERY_COUNT} variations")

            st.caption("**Post-Processing:**")
            st.caption(f"• Reranking: {'✅' if CONFIG.ENABLE_RERANKING else '❌'}")
            st.caption(
                f"• MMR (λ={CONFIG.MMR_LAMBDA}): {'✅' if CONFIG.ENABLE_MMR else '❌'}"
            )
            st.caption(f"• CRAG: {'✅' if CONFIG.ENABLE_CRAG else '❌'}")
            st.caption(
                f"• Context Compression: {'✅' if CONFIG.ENABLE_CONTEXT_COMPRESSION else '❌'}"
            )

            st.caption("**Caching:**")
            st.caption(
                f"• Semantic Cache: {'✅' if CONFIG.ENABLE_SEMANTIC_CACHE else '❌'}"
            )

            st.caption("**Arabic:**")
            st.caption(
                f"• Morphology: {'✅' if CONFIG.ENABLE_ARABIC_MORPHOLOGY else '❌'}"
            )

        with st.expander("🔬 Run Evaluation"):
            if st.button("Run Sample Evaluation", use_container_width=True):
                if not st.session_state.api_key:
                    st.error("API key required")
                else:
                    with st.spinner("Running evaluation..."):
                        llm_client = LLMClient(
                            st.session_state.selected_provider, st.session_state.api_key
                        )
                        embedder = load_embedder()
                        _, collection = get_chroma_collection()

                        pipeline = RAGPipeline(llm_client, embedder, collection)
                        evaluator = RAGEvaluator(pipeline, llm_client)

                        eval_queries = create_eval_dataset_from_corpus(collection, llm_client, n=5)
                        if not eval_queries:
                            st.error("No documents found (or could not generate eval set). Ingest documents first.")
                            st.stop()

                        results = evaluator.evaluate_dataset(eval_queries)

                        # Save results
                        with open(CONFIG.EVAL_RESULTS_PATH, "w") as f:
                            json.dump(results, f, indent=2)

                        st.success("Evaluation complete!")
                        UIRenderer.render_eval_metrics(results.get("overall", {}))

        st.markdown("---")
        col1, col2 = st.columns(2)

        with col1:
            if st.button("🗑️ Clear Chat", use_container_width=True):
                st.session_state.chat_history = []
                st.rerun()

        with col2:
            if st.button("🔄 Reset DB", use_container_width=True):
                try:
                    client, _ = get_chroma_collection()
                    client.reset()
                    for key in ["chroma_client", "chroma_collection", "bm25_manager"]:
                        if key in st.session_state:
                            del st.session_state[key]
                    st.session_state.semantic_cache = OrderedDict()

                    # Delete BM25 index file
                    bm25_path = Path(CONFIG.BM25_INDEX_PATH)
                    if bm25_path.exists():
                        bm25_path.unlink()

                    st.success("Database reset!")
                    st.rerun()
                except Exception as e:
                    st.error(str(e))

        if CONFIG.ENABLE_SEMANTIC_CACHE:
            if st.button("🧹 Clear Cache", use_container_width=True):
                st.session_state.semantic_cache = OrderedDict()
                st.success("Cache cleared!")
                st.rerun()


# ============================================================================
# MAIN APPLICATION
# ============================================================================
def main():
    """Main application entry point."""
    render_sidebar()

    # Welcome header when no chat history
    if not st.session_state.chat_history:
        st.markdown(
            """
        <div class="app-header">
            <div class="app-title">🎓 University Tutor - AI Engineering Specialist</div>
            <div class="app-subtitle">Ask questions about your lecture materials in English or Arabic</div>
        </div>
        """,
            unsafe_allow_html=True,
        )

        provider_name = LLM_PROVIDERS[st.session_state.selected_provider]["name"]
        model_name = LLM_PROVIDERS[st.session_state.selected_provider]["model"]
        st.markdown(
            f"""
        <div style="text-align: center; margin-bottom: 2rem;">
            <span class="provider-badge">⚡ Powered by {provider_name} • {model_name.split('/')[-1]}</span>
            <br><br>
            <span class="provider-badge">🚀 Enhanced RAG: True Hybrid Search (BM25+Vector) • HyDE • 
            Arabic Morphology • MMR Diversity • Context Compression • CRAG</span>
        </div>
        """,
            unsafe_allow_html=True,
        )

    # Check for API key
    if not st.session_state.api_key:
        st.info(
            "👈 Select a provider and enter your API key in the sidebar to get started"
        )
        st.stop()

    # Initialize LLM client
    llm_client = LLMClient(st.session_state.selected_provider, st.session_state.api_key)

    # Render chat history
    followup_clicked = None
    for msg_idx, msg in enumerate(st.session_state.chat_history):
        result = UIRenderer.render_message(
            msg["role"],
            msg["content"],
            msg.get("sources"),
            msg.get("followups"),
            msg.get("question_lang", "en"),
            message_idx=msg_idx,
        )
        if result:
            followup_clicked = result

    # Handle followup clicks
    if followup_clicked:
        st.session_state.pending_question = followup_clicked
        st.rerun()

    # Chat input
    question = st.chat_input(
        "Ask anything about your lecture materials... / اسأل أي سؤال عن محاضراتك..."
    )

    # Handle pending questions from followups
    if "pending_question" in st.session_state:
        question = st.session_state.pending_question
        del st.session_state.pending_question

    # Process question
    if question:
        question = question.strip()
        if question:
            question_lang = LanguageUtils.detect_language(question)
            st.session_state.chat_history.append({"role": "user", "content": question})
            UIRenderer.render_message("user", question)

            with st.spinner("🔍 Searching lecture materials and analyzing..."):
                embedder = load_embedder()
                _, collection = get_chroma_collection()

                pipeline = RAGPipeline(llm_client, embedder, collection)
                answer, sources, followups, debug = pipeline.answer(question)

                # Show debug info
                debug_parts = []
                if debug.get("cache_hit"):
                    sim = debug.get("cache_similarity")
                    debug_parts.append(f"Cache hit ({sim:.3f})" if sim else "Cache hit")
                if debug.get("comparison_detected"):
                    topics = debug.get("comparison_topics", [])
                    debug_parts.append(f"Comparison: {' vs '.join(topics)}")
                if debug.get("crag_used"):
                    m = debug.get("crag_metrics", {})
                    debug_parts.append(
                        f"CRAG (max_rel={m.get('max_relevance', 0):.2f})"
                    )
                if debug.get("hyde_used"):
                    debug_parts.append("HyDE")
                if debug.get("bm25_used"):
                    debug_parts.append("BM25")
                if debug.get("mmr_used"):
                    debug_parts.append("MMR")

                if debug_parts:
                    st.caption(f"🔧 Features used: {' • '.join(debug_parts)}")

            st.session_state.chat_history.append(
                {
                    "role": "assistant",
                    "content": answer,
                    "sources": sources,
                    "followups": followups,
                    "question_lang": question_lang,
                }
            )
            st.rerun()

    # Footer
    st.markdown("---")
    provider_name = LLM_PROVIDERS[st.session_state.selected_provider]["name"]
    st.markdown(
        f"""
    <div style="text-align: center; color: #666; font-size: 0.85rem;">
        <strong>Enhanced RAG Architecture:</strong> True Hybrid Search (BM25 + Vector) • 
        HyDE • Arabic Morphology (CAMeL) • Multi-Query • Cross-lingual • 
        Reranking • MMR Diversity • CRAG • Context Compression • Semantic Cache<br>
        <strong>Powered by:</strong> {provider_name} • ChromaDB • E5-multilingual-small
    </div>
    """,
        unsafe_allow_html=True,
    )


if __name__ == "__main__":
    main()