import streamlit as st
from typing import List, Dict, Tuple
from pptx import Presentation
from pypdf import PdfReader
from config import Config
from languageUtils import LanguageUtils

class DocumentProcessor:
    """Document extraction and chunking."""

    @staticmethod
    def extract_from_pdf(file) -> List[Tuple[int, str]]:
        """Extract text from PDF."""
        try:
            reader = PdfReader(file)
            pages = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append((i, text.strip()))
            return pages
        except Exception as e:
            st.error(f"PDF error: {e}")
            return []

    @staticmethod
    def extract_from_ppt(file) -> List[Tuple[int, str]]:
        """Extract text from PowerPoint."""
        try:
            prs = Presentation(file)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append((i, "\n".join(texts)))
            return slides
        except Exception as e:
            st.error(f"PPT error: {e}")
            return []

    @staticmethod
    def chunk_text(text: str, chunk_size: int, overlap: int) -> List[str]:
        """Split text into overlapping chunks."""
        if len(text) <= chunk_size:
            return [text]

        chunks = []
        start = 0

        while start < len(text):
            end = start + chunk_size

            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings
                for sep in [". ", ".\n", "؟ ", "。", "\n\n"]:
                    last_sep = text[start:end].rfind(sep)
                    if last_sep > chunk_size // 2:
                        end = start + last_sep + len(sep)
                        break

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap
            if start >= len(text):
                break

        return chunks

    @classmethod
    def process_document(cls, file) -> List[Dict]:
        """Process a document into chunks with metadata."""
        filename = file.name.lower()

        if filename.endswith(".pdf"):
            pages = cls.extract_from_pdf(file)
            doc_type = "pdf"
        elif filename.endswith((".ppt", ".pptx")):
            pages = cls.extract_from_ppt(file)
            doc_type = "ppt"
        else:
            return []

        all_chunks = []

        for page_num, text in pages:
            lang = LanguageUtils.detect_language(text)

            # Create parent chunks
            parent_chunks = cls.chunk_text(
                text, Config.PARENT_CHUNK_SIZE, Config.PARENT_OVERLAP
            )

            for p_idx, parent_text in enumerate(parent_chunks):
                # Create child chunks from parent
                child_chunks = cls.chunk_text(
                    parent_text, Config.CHILD_CHUNK_SIZE, Config.CHILD_OVERLAP
                )

                for c_idx, child_text in enumerate(child_chunks):
                    all_chunks.append(
                        {
                            "child_text": child_text,
                            "parent_text": parent_text,
                            "metadata": {
                                "source": file.name,
                                "page": page_num,
                                "type": doc_type,
                                "lang": lang,
                                "parent_idx": p_idx,
                                "child_idx": c_idx,
                            },
                        }
                    )

        return all_chunks